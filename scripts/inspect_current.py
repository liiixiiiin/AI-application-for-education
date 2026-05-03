"""Quick inspection of the current 中期答辩PPT.pptx — list all shape types
on each slide and flag which slides contain user-added pictures."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

PPTX = Path("/Users/daniel/2026design/App/中期答辩PPT.pptx")


def in_(emu_val) -> float:
    return emu_val / 914400.0


def main() -> None:
    prs = Presentation(PPTX)
    print(f"Slides: {len(prs.slides)}")
    for i, sld in enumerate(prs.slides, 1):
        pics = []
        boxes = 0
        autos = 0
        groups = 0
        others = 0
        for shp in sld.shapes:
            st = shp.shape_type
            try:
                w = in_(shp.width); h = in_(shp.height)
                top = in_(shp.top); left = in_(shp.left)
            except Exception:
                w = h = top = left = -1
            if st == 13:
                pics.append((w, h, top, left, getattr(shp, "name", "")))
            elif st == 17:
                boxes += 1
            elif st == 1:
                autos += 1
            elif st == 6:
                groups += 1
            else:
                others += 1
        n_logo = sum(1 for w, h, t, l, _ in pics if t < 1.0 and l > 9.5 and w < 3.0)
        n_user = len(pics) - n_logo
        print(f"\n=== Slide {i} === text:{boxes} auto:{autos} groups:{groups} pics:{len(pics)} (logos:{n_logo} user/large:{n_user}) others:{others}")
        for w, h, t, l, name in pics:
            tag = "logo" if (t < 1.0 and l > 9.5 and w < 3.0) else "USER/LARGE"
            print(f"   pic [{tag:10s}] w={w:5.2f} h={h:5.2f} top={t:5.2f} left={l:5.2f}  {name}")


if __name__ == "__main__":
    main()
