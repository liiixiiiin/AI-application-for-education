"""Inspect shapes in masters/layouts/slides to find decoration sources."""
from pathlib import Path
from pptx import Presentation

TEMPLATE_PATH = Path(
    "/Users/daniel/答辩PPT/中国石油大学（华东）PPT模板/学海遨游，逐梦石大-学术系列_副本.pptx"
)


def describe_shape(s) -> str:
    typ = s.shape_type
    nm = s.name
    has_text = ""
    try:
        if s.has_text_frame and s.text_frame.text:
            t = s.text_frame.text.replace("\n", " | ")[:40]
            has_text = f' text={t!r}'
    except Exception:
        pass
    has_pic = ""
    try:
        if s.shape_type == 13:  # picture
            has_pic = " [PICTURE]"
    except Exception:
        pass
    return f"type={typ} name={nm!r}{has_text}{has_pic}"


def main() -> None:
    prs = Presentation(TEMPLATE_PATH)
    print("=== Master 0 ===")
    m = prs.slide_masters[0]
    for s in m.shapes:
        print(f"  M0: {describe_shape(s)}")
    print("\n=== Master 0 / Layout 0 (空白) ===")
    for s in m.slide_layouts[0].shapes:
        print(f"  L0: {describe_shape(s)}")
    print("\n=== Master 0 / Layout 9 (标题和文本) ===")
    for s in m.slide_layouts[9].shapes:
        print(f"  L9: {describe_shape(s)}")
    print("\n=== Master 0 / Layout 10 (标题幻灯片) ===")
    for s in m.slide_layouts[10].shapes:
        print(f"  L10: {describe_shape(s)}")

    print("\n=== Slide 1 (cover) - shape count:", len(prs.slides[0].shapes), "===")
    for s in prs.slides[0].shapes:
        print(f"  S1: {describe_shape(s)}")
    print(f"  -> layout: {prs.slides[0].slide_layout.name!r}")
    print(f"  -> master: {prs.slides[0].slide_layout.slide_master.name!r}")

    print("\n=== Slide 2 (content) - shape count:", len(prs.slides[1].shapes), "===")
    for s in prs.slides[1].shapes:
        print(f"  S2: {describe_shape(s)}")
    print(f"  -> layout: {prs.slides[1].slide_layout.name!r}")

    print("\n=== Slide 27 (thanks) ===")
    for s in prs.slides[26].shapes:
        print(f"  S27: {describe_shape(s)}")
    print(f"  -> layout: {prs.slides[26].slide_layout.name!r}")


if __name__ == "__main__":
    main()
