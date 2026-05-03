"""Incremental update of /Users/daniel/2026design/App/中期答辩PPT.pptx.

This script does NOT regenerate the deck from scratch. It edits the user's
current PPT in-place, preserving every picture the user has manually added
(截图、ER 图、政策图 等). Only the pages explicitly listed below are
re-built; everything else is left intact except for two global passes:

    1. font_size  +1 pt  (everything from 8 pt to 14 pt grows by 1 pt)
    2. fancy emoji → simple icons / removed

Run:
    .ppt_venv/bin/python scripts/update_defense_pptx.py
"""

from __future__ import annotations

import io
import re
import sys
from pathlib import Path

# Re-use the helpers we already wrote in build_defense_pptx.py
sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_defense_pptx import (  # noqa: E402  (path manipulation above)
    COLOR_ACCENT,
    COLOR_GRAY_BG,
    COLOR_LIGHT_BG,
    COLOR_PRIMARY,
    COLOR_SUBTLE,
    COLOR_TEXT,
    EN_FONT,
    SLIDE_W_IN,
    add_filled_rect,
    add_image_placeholder,
    add_paragraphs,
    add_textbox,
    add_title_bar,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PPT_PATH = Path("/Users/daniel/2026design/App/中期答辩PPT.pptx")

# =============================================================================
# Common helpers
# =============================================================================


def _in(emu_val) -> float:
    return emu_val / 914400.0


def is_decoration_logo(shape) -> bool:
    """The little school logo at the top-right corner is a decoration."""
    if shape.shape_type != 13:  # PICTURE
        return False
    try:
        w = _in(shape.width)
        top = _in(shape.top)
        left = _in(shape.left)
    except Exception:
        return False
    # Top-right small logo (matches the template's stamp)
    return top < 1.0 and left > 9.5 and w < 3.0


def is_user_picture(shape) -> bool:
    """A picture the user manually added (large content image)."""
    if shape.shape_type != 13:
        return False
    if is_decoration_logo(shape):
        return False
    return True


def clear_slide_keep_user_pictures(slide) -> None:
    """Wipe everything that is *not* a picture; keep all PICTURE shapes
    (both the small logo and any large user-added images)."""
    to_remove = []
    for shp in slide.shapes:
        if shp.shape_type == 13:
            continue  # keep all pictures
        to_remove.append(shp)
    for shp in to_remove:
        sp = shp._element
        sp.getparent().remove(sp)


def clear_slide_keep_logo_only(slide) -> None:
    """Remove every shape except the top-right small school logo."""
    to_remove = []
    for shp in slide.shapes:
        if is_decoration_logo(shp):
            continue
        to_remove.append(shp)
    for shp in to_remove:
        sp = shp._element
        sp.getparent().remove(sp)


# =============================================================================
# Tiny icon primitives (replace emoji)
# =============================================================================


def add_icon_circle(slide, left: float, top: float, size: float,
                    fill: RGBColor, text: str = "",
                    text_color: RGBColor | None = None) -> None:
    """A solid coloured circle, optionally with a 1-2 char label."""
    sh = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(size), Inches(size),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    if text:
        tf = sh.text_frame
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        from pptx.enum.text import MSO_ANCHOR as _A
        tf.vertical_anchor = _A.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(int(size * 9))
        run.font.bold = True
        run.font.color.rgb = text_color or RGBColor(0xFF, 0xFF, 0xFF)


def add_icon_square(slide, left: float, top: float, size: float,
                    fill: RGBColor) -> None:
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(size), Inches(size),
    )
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()


# =============================================================================
# 1. Global passes: font size + emoji cleanup
# =============================================================================


# Replacement table: fancy emoji → simple ASCII / Unicode glyph
EMOJI_MAP = {
    "💡": "★",
    "⭐": "★",
    "❗": "!",
    "❓": "?",
    "✦": "◆",
    "✏️": "▎",
    "✏": "▎",
    "📷": "■",
    "📥": "▼",
    "📋": "■",
    "📊": "■",
    "📈": "■",
    "📐": "■",
    "📚": "■",
    "📝": "■",
    "📁": "■",
    "📂": "■",
    "🔧": "●",
    "🛡️": "●",
    "🛡": "●",
    "🎯": "●",
    "🎓": "●",
    "🚀": "●",
    "🔄": "↻",
    "🌐": "●",
    "🧭": "●",
    "🧩": "●",
    "🔥": "●",
    "💬": "■",
    "🗂": "■",
    "🗂️": "■",
    # variation selectors / ZWJ artifacts can be removed:
    "\ufe0f": "",
    "\u200d": "",
}


def clean_emojis_in_text(text: str) -> str:
    out = text
    for k, v in EMOJI_MAP.items():
        out = out.replace(k, v)
    return out


_PAGE_NUM_RE = re.compile(r"(\d+)\s*/\s*13\b")


def normalize_paragraph(paragraph, font_delta_pt: int = 1) -> None:
    """Bump every run's font size by `font_delta_pt`, replace fancy
    emojis with simple icons, and fix outdated "X/13" page numbers."""
    for run in paragraph.runs:
        # Emoji cleanup + page-number fix (now 12 pages)
        if run.text:
            new_t = clean_emojis_in_text(run.text)
            new_t = _PAGE_NUM_RE.sub(r"\1 / 12", new_t)
            if new_t != run.text:
                run.text = new_t
        # Font size bump
        try:
            sz = run.font.size
        except Exception:
            sz = None
        if sz is not None:
            try:
                run.font.size = Pt(sz.pt + font_delta_pt)
            except Exception:
                pass


def normalize_slide(slide, font_delta_pt: int = 1) -> None:
    for shp in slide.shapes:
        if shp.has_text_frame:
            for p in shp.text_frame.paragraphs:
                normalize_paragraph(p, font_delta_pt)


# =============================================================================
# 2. Slide rebuilders
# =============================================================================


def add_subtle_card(slide, left, top, width, height,
                    fill=COLOR_LIGHT_BG, corner=0.06):
    return add_filled_rect(slide, left, top, width, height,
                           fill=fill, corner=corner)


def rebuild_p3_objective(slide) -> None:
    """Page 3 — 研究目标. Pain → Solution comparison + 4 small image
    placeholders (one per solution) for richer layout."""
    clear_slide_keep_logo_only(slide)
    add_title_bar(slide, "研究目标 — 针对痛点的系统性对策")

    # Top banner
    add_filled_rect(slide, 0.55, 1.20, 12.30, 0.78, fill=COLOR_PRIMARY)
    add_textbox(
        slide,
        left=0.85, top=1.20, width=11.70, height=0.78,
        text="构建教师 / 学生 / 管理三端一体的 AI 教学辅助系统，"
             "针对四大痛点提供 RAG · 知识追踪 · 个性化推荐三位一体的系统性对策",
        size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    pain_color = RGBColor(0xC0, 0x39, 0x2B)
    solu_color = RGBColor(0x2E, 0x9C, 0x6E)

    # Side labels
    add_textbox(slide, left=0.55, top=2.20, width=2.0, height=0.32,
                text="痛点", size=15, bold=True, color=pain_color)
    add_textbox(slide, left=0.55, top=4.40, width=2.0, height=0.32,
                text="对策", size=15, bold=True, color=solu_color)
    add_textbox(slide, left=0.55, top=5.95, width=2.0, height=0.32,
                text="系统体现", size=15, bold=True, color=COLOR_PRIMARY)

    pairs = [
        {
            "pain_title": "备课负担重",
            "pain_desc": "章节素材整理、题目设计耗时数小时",
            "solu_title": "教师备课辅助",
            "solu_desc": "章节讲解提纲自动生成\n含教学目标 / 重难点 / 课堂流程",
            "demo_label": "备课提纲生成结果",
        },
        {
            "pain_title": "题库更新慢",
            "pain_desc": "考核内容陈旧、与课程内容脱节",
            "solu_title": "智能出题",
            "solu_desc": "按知识点定向生成 4 种题型\n单选 / 判断 / 填空 / 简答",
            "demo_label": "出题结果列表",
        },
        {
            "pain_title": "个性化缺位",
            "pain_desc": "一对多教学难以兼顾差异化",
            "solu_title": "知识追踪 + 推荐",
            "solu_desc": "EMA 动态更新掌握度\n推送针对薄弱知识点的练习",
            "demo_label": "个性化推荐",
        },
        {
            "pain_title": "数据成孤岛",
            "pain_desc": "教与学过程数据分散、难以闭环",
            "solu_title": "数据闭环",
            "solu_desc": "练习 → 评估 → EMA 更新\n作答数据自动驱动推荐",
            "demo_label": "学情数据闭环",
        },
    ]

    col_w = 2.85
    col_gap = 0.22
    cols_total_w = 4 * col_w + 3 * col_gap
    start_left = 0.55 + (12.30 - cols_total_w) / 2

    pain_top, pain_h = 2.55, 1.55
    arrow_top, arrow_h = 4.14, 0.20
    solu_top, solu_h = 4.74, 1.18
    img_top, img_h = 5.96, 1.20

    for i, p in enumerate(pairs):
        x = start_left + i * (col_w + col_gap)

        # Pain card
        add_filled_rect(slide, x, pain_top, col_w, pain_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, x, pain_top, col_w, 0.50,
                        fill=pain_color, corner=0.06)
        # icon (small white square) instead of emoji
        add_icon_square(slide, x + 0.18, pain_top + 0.13, 0.24,
                        RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left=x + 0.55, top=pain_top + 0.05,
                    width=col_w - 0.65, height=0.40,
                    text=p["pain_title"], size=15, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=x + 0.15, top=pain_top + 0.60,
                    width=col_w - 0.30, height=pain_h - 0.70,
                    text=p["pain_desc"], size=13, color=COLOR_TEXT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.4)

        # Down arrow
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(x + col_w / 2 - 0.18), Inches(arrow_top),
            Inches(0.36), Inches(arrow_h),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_PRIMARY
        arrow.line.fill.background()

        # Solution card
        add_filled_rect(slide, x, solu_top, col_w, solu_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, x, solu_top, col_w, 0.50,
                        fill=solu_color, corner=0.06)
        add_icon_square(slide, x + 0.18, solu_top + 0.13, 0.24,
                        RGBColor(0xFF, 0xFF, 0xFF))
        add_textbox(slide, left=x + 0.55, top=solu_top + 0.05,
                    width=col_w - 0.65, height=0.40,
                    text=p["solu_title"], size=15, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=x + 0.15, top=solu_top + 0.60,
                    width=col_w - 0.30, height=solu_h - 0.70,
                    text=p["solu_desc"], size=12, color=COLOR_TEXT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    line_spacing=1.35)

        # Image placeholder underneath each solution
        add_image_placeholder(
            slide, left=x, top=img_top, width=col_w, height=img_h,
            label=p["demo_label"],
        )

    # Bottom banner
    add_filled_rect(slide, 0.55, 7.20, 11.20, 0.30, fill=COLOR_PRIMARY,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE, corner=0.40)
    add_textbox(
        slide, left=0.55, top=7.20, width=11.20, height=0.30,
        text="痛点逐一被对策接住 — 形成可量化、可落地、可持续演进的智能教学闭环",
        size=12, bold=True, color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


def rebuild_p5_architecture(slide) -> None:
    """Page 5 — 综合架构 & 技术栈. Each tier hosts a row of small "blocks"
    (rounded rectangles) to make it more visual."""
    # Preserve any user-added picture (e.g. ER diagram if they pasted one
    # here). We still strip the small logo and re-add later? No — keep
    # all pictures including logo:
    clear_slide_keep_user_pictures(slide)
    add_title_bar(slide, "系统总体架构与技术栈")

    add_textbox(
        slide, left=0.55, top=1.05, width=12.30, height=0.30,
        text="◆ 6 层模块化架构  ·  前后端分离  ·  模型可替换  ·  流式优先  ·  数据闭环",
        size=13, bold=True, color=COLOR_PRIMARY,
    )

    tiers = [
        {
            "name": "用户层", "name_en": "User",
            "color": RGBColor(0x37, 0x4F, 0x7E),
            "blocks": [
                ("教师", "Teacher"),
                ("学生", "Student"),
                ("管理员", "Admin"),
                ("浏览器", "Chrome / Edge"),
                ("RBAC 权限", "三角色控制"),
            ],
        },
        {
            "name": "展示层", "name_en": "Presentation",
            "color": RGBColor(0x4A, 0x90, 0xE2),
            "blocks": [
                ("Vue 3", "组合式 API"),
                ("Vite", "极速构建"),
                ("Element Plus", "UI 组件库"),
                ("Lucide", "图标"),
                ("Pinia", "状态管理"),
                ("ECharts", "图表"),
            ],
        },
        {
            "name": "业务层", "name_en": "Business",
            "color": RGBColor(0x2E, 0x9C, 0x6E),
            "blocks": [
                ("auth", "登录/会话"),
                ("courses", "课程"),
                ("knowledge_base", "知识库"),
                ("rag_qa", "RAG 问答"),
                ("exercises", "智能出题"),
                ("lesson_plans", "备课"),
                ("knowledge_tracking", "学情"),
                ("conversations", "多轮对话"),
            ],
        },
        {
            "name": "编排层", "name_en": "Orchestration",
            "color": RGBColor(0xE5, 0xA5, 0x3B),
            "blocks": [
                ("LangChain", "Pipeline"),
                ("Prompt 工程", "Few-shot"),
                ("JSON Schema", "结构化输出"),
                ("RAG Pipeline", "混合检索"),
                ("EMA 追踪", "掌握度"),
            ],
        },
        {
            "name": "AI 服务层", "name_en": "AI Services",
            "color": RGBColor(0xC0, 0x55, 0x55),
            "blocks": [
                ("Qwen-Plus", "对话生成"),
                ("Qwen-Turbo", "快速响应"),
                ("text-embedding-v3", "向量化"),
                ("gte-rerank", "重排序"),
                ("RAGAS", "评测"),
            ],
        },
        {
            "name": "数据层", "name_en": "Data",
            "color": RGBColor(0x69, 0x40, 0x90),
            "blocks": [
                ("SQLite", "8 张业务表"),
                ("ChromaDB", "向量索引 HNSW"),
                ("BM25 + jieba", "倒排索引"),
                ("本地文件", "切片 / 原文"),
                ("JSON Schema", "对话引用"),
            ],
        },
    ]

    base_top = 1.45
    tier_h = 0.86
    gap = 0.05
    layer_w = 12.30
    label_w = 1.30

    for ti, t in enumerate(tiers):
        top = base_top + ti * (tier_h + gap)

        # Tier background card
        add_filled_rect(slide, 0.55, top, layer_w, tier_h,
                        fill=COLOR_LIGHT_BG, corner=0.04)
        # Left colour label
        add_filled_rect(slide, 0.55, top, label_w, tier_h,
                        fill=t["color"], corner=0.04)
        add_textbox(
            slide, left=0.55, top=top + 0.07, width=label_w, height=0.42,
            text=t["name"], size=15, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left=0.55, top=top + 0.50, width=label_w, height=0.30,
            text=t["name_en"], size=9,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, font_name=EN_FONT,
        )

        # Blocks: stretch evenly across remaining width
        n = len(t["blocks"])
        avail_w = layer_w - label_w - 0.20  # 0.10 padding each side
        block_gap = 0.10
        block_w = (avail_w - (n - 1) * block_gap) / n
        block_h = tier_h - 0.16
        bx0 = 0.55 + label_w + 0.10
        by = top + 0.08
        for bi, (b_title, b_sub) in enumerate(t["blocks"]):
            bx = bx0 + bi * (block_w + block_gap)
            # Block body (white) with coloured top bar
            add_filled_rect(slide, bx, by, block_w, block_h,
                            fill=RGBColor(0xFF, 0xFF, 0xFF), corner=0.06)
            add_filled_rect(slide, bx, by, block_w, 0.16,
                            fill=t["color"], corner=0.06)
            add_textbox(
                slide, left=bx + 0.04, top=by + 0.18,
                width=block_w - 0.08, height=0.30,
                text=b_title, size=11, bold=True, color=t["color"],
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )
            add_textbox(
                slide, left=bx + 0.04, top=by + 0.46,
                width=block_w - 0.08, height=0.20,
                text=b_sub, size=9, color=COLOR_TEXT,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

    # Bottom flow note
    add_textbox(
        slide, left=0.55, top=7.10, width=12.30, height=0.32,
        text="数据闭环：学生作答 → exercise_attempts → EMA 更新 knowledge_mastery → 识别薄弱点 → 个性化推荐",
        size=12, bold=True, color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )


def rebuild_p6_database(slide) -> None:
    """Page 6 — 数据库设计 with tricky technical highlights.

    Layout:
        - top: 8 tables grid (compact)
        - middle-left: ER 占位/图（保留用户已加图）
        - middle-right: 5 个 Tricky Tech 卡片
        - bottom: 紧凑的数据闭环 hint
    """
    # Keep user pictures (incl. ER diagram if they pasted one)
    clear_slide_keep_user_pictures(slide)
    add_title_bar(slide, "核心数据库设计与技术亮点")

    # ---------- 8 tables grid (top, compact) ----------
    add_textbox(
        slide, left=0.55, top=1.05, width=12.30, height=0.30,
        text="◆ 8 张核心业务表 — 按 基础数据 / 对话 / 知识追踪 三组着色",
        size=13, bold=True, color=COLOR_PRIMARY,
    )

    tables = [
        ("users", "用户账号 三角色", COLOR_ACCENT),
        ("sessions", "登录会话 token", COLOR_ACCENT),
        ("courses", "课程信息", COLOR_ACCENT),
        ("knowledge_points", "课程知识点", COLOR_ACCENT),
        ("conversations", "对话会话", RGBColor(0x2E, 0x9C, 0x6E)),
        ("messages", "对话消息 含引用", RGBColor(0x2E, 0x9C, 0x6E)),
        ("knowledge_mastery", "EMA 掌握度 ★", RGBColor(0xD9, 0x82, 0x2E)),
        ("exercise_attempts", "作答记录 ★", RGBColor(0xD9, 0x82, 0x2E)),
    ]

    cell_w = 2.90
    cell_h = 0.85
    gap_x = 0.20
    gap_y = 0.12
    start_left = (SLIDE_W_IN - 4 * cell_w - 3 * gap_x) / 2
    base_top = 1.40
    for idx, (name, desc, color) in enumerate(tables):
        col = idx % 4
        row = idx // 4
        x = start_left + col * (cell_w + gap_x)
        y = base_top + row * (cell_h + gap_y)
        add_filled_rect(slide, x, y, cell_w, cell_h, fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, x, y, cell_w, 0.32, fill=color, corner=0.06)
        add_textbox(slide, left=x, top=y + 0.02, width=cell_w, height=0.30,
                    text=name, size=12, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                    font_name=EN_FONT)
        add_textbox(slide, left=x + 0.08, top=y + 0.34,
                    width=cell_w - 0.16, height=cell_h - 0.36,
                    text=desc, size=11, color=COLOR_TEXT,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---------- ER area (left) + Tricky Tech cards (right) ----------
    section_top = 3.50
    section_h = 3.40

    # Left: ER label only (the actual picture is preserved from user paste)
    er_left = 0.55
    er_w = 5.40
    add_filled_rect(slide, er_left, section_top, er_w, 0.36,
                    fill=COLOR_PRIMARY)
    add_textbox(
        slide, left=er_left, top=section_top, width=er_w, height=0.36,
        text="◆ ER 关系示意图（自动生成自 db.py）",
        size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # If user did NOT paste a picture into this slot, draw a hint.
    has_user_pic_in_er_slot = any(
        is_user_picture(s) and _in(s.left) < er_left + er_w
        and _in(s.top) > section_top - 0.2
        for s in slide.shapes
    )
    if not has_user_pic_in_er_slot:
        add_filled_rect(
            slide, er_left, section_top + 0.36, er_w, section_h - 0.36,
            fill=RGBColor(0xFA, 0xFC, 0xFE), corner=0.04,
        )
        add_textbox(
            slide, left=er_left + 0.20, top=section_top + 0.55,
            width=er_w - 0.40, height=section_h - 0.70,
            text=(
                "ER 图源文件已生成：\n"
                "data/er-diagram/er.svg  (矢量)\n"
                "data/er-diagram/er.png  (位图)\n\n"
                "拖入此处即可"
            ),
            size=12, color=RGBColor(0x66, 0x66, 0x66),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.35,
        )

    # Right: 5 Tricky tech highlight cards (vertical stack)
    tt_left = 6.20
    tt_w = 6.65
    add_filled_rect(slide, tt_left, section_top, tt_w, 0.36,
                    fill=RGBColor(0xD9, 0x82, 0x2E))
    add_textbox(
        slide, left=tt_left, top=section_top, width=tt_w, height=0.36,
        text="◆ 数据库设计的 5 个技术亮点",
        size=12, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    tricky = [
        ("HNSW",
         "ChromaDB 默认 HNSW 向量索引 — 高维近似最近邻 O(log n)，"
         "千万级片段下毫秒级召回"),
        ("UPSERT",
         "knowledge_mastery 复合 UNIQUE (student, course, point) — "
         "原子 upsert 保证掌握度记录唯一不重复"),
        ("EMA",
         "α=0.7 指数平滑：mastery_new = α·old + (1-α)·score — "
         "兼顾历史与最新作答，平滑抗抖动"),
        ("JSON-in-TEXT",
         "messages.citations / exercise_attempts.knowledge_points — "
         "TEXT 存 JSON 数组，schema 灵活、查询时按需解码"),
        ("FTS5 预留",
         "knowledge_points 表预留主键索引 — 后续可一键挂载 SQLite "
         "FTS5 实现关键词全文检索（计划升级）"),
    ]
    card_h = (section_h - 0.36 - (len(tricky) - 1) * 0.06) / len(tricky)
    card_top0 = section_top + 0.42
    for i, (tag, desc) in enumerate(tricky):
        cy = card_top0 + i * (card_h + 0.06)
        add_filled_rect(slide, tt_left, cy, tt_w, card_h,
                        fill=COLOR_LIGHT_BG, corner=0.04)
        # Tag chip on left
        chip_w = 1.45
        add_filled_rect(slide, tt_left, cy, chip_w, card_h,
                        fill=COLOR_PRIMARY, corner=0.04)
        add_textbox(
            slide, left=tt_left, top=cy, width=chip_w, height=card_h,
            text=tag, size=11, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            font_name=EN_FONT,
        )
        add_textbox(
            slide, left=tt_left + chip_w + 0.10, top=cy + 0.05,
            width=tt_w - chip_w - 0.20, height=card_h - 0.10,
            text=desc, size=10, color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.30,
        )

    # ---------- Bottom flow note ----------
    add_textbox(
        slide, left=0.55, top=7.05, width=12.30, height=0.32,
        text="数据闭环：作答 → exercise_attempts → EMA 平滑更新 knowledge_mastery → 识别薄弱知识点 → 个性化推荐 → 推送学生",
        size=12, bold=True, color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )


def rebuild_p9_lesson_plan(slide) -> None:
    """Page 9 — 教师备课辅助.

    The user has already pasted ~7 screenshots in the bottom strip
    (top ≥ 5.69"), so we keep that area FREE and confine all newly added
    text/cards to the top 5.5". This naturally gives the user's
    screenshots the dominant share of the page.
    """
    clear_slide_keep_user_pictures(slide)
    add_title_bar(slide, "模块 2 — 教师备课辅助（章节讲解提纲）")

    # ---------- Left narrow column (3.10") ----------
    left_w = 3.10

    # Input params card  (1.20 – 3.10)
    add_filled_rect(slide, 0.55, 1.15, left_w, 0.40, fill=COLOR_PRIMARY,
                    corner=0.06)
    add_textbox(
        slide, left=0.55, top=1.15, width=left_w, height=0.40,
        text="输入参数",
        size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_filled_rect(slide, 0.55, 1.55, left_w, 1.55, fill=COLOR_LIGHT_BG,
                    corner=0.06)
    box = slide.shapes.add_textbox(
        Inches(0.70), Inches(1.62), Inches(left_w - 0.30), Inches(1.45)
    )
    add_paragraphs(
        box.text_frame,
        [
            ("课程", {"size": 11, "bold": True, "color": COLOR_PRIMARY,
                      "space_after": 1}),
            ("  大模型应用开发",
             {"size": 11, "color": COLOR_TEXT, "space_after": 4}),
            ("章节", {"size": 11, "bold": True, "color": COLOR_PRIMARY,
                      "space_after": 1}),
            ("  第 3 章 RAG 检索增强生成",
             {"size": 11, "color": COLOR_TEXT, "space_after": 4}),
            ("课时 / 知识点",
             {"size": 11, "bold": True, "color": COLOR_PRIMARY,
              "space_after": 1}),
            ("  90 分钟 · 向量检索 / BM25 / Rerank",
             {"size": 11, "color": COLOR_TEXT}),
        ],
    )

    # Generation flow card  (3.20 – 5.30)
    add_filled_rect(slide, 0.55, 3.20, left_w, 0.40, fill=COLOR_ACCENT,
                    corner=0.06)
    add_textbox(
        slide, left=0.55, top=3.20, width=left_w, height=0.40,
        text="生成流程", size=14, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_filled_rect(slide, 0.55, 3.60, left_w, 1.70, fill=COLOR_LIGHT_BG,
                    corner=0.06)
    flow_items = [
        "① 检索知识库 Top-N 片段",
        "② 拼装 Prompt + Few-shot",
        "③ Qwen-Plus 调用",
        "④ 6 大模块结构化输出",
        "⑤ 教师审核 / 微调 / 导出",
    ]
    fbox = slide.shapes.add_textbox(
        Inches(0.70), Inches(3.68), Inches(left_w - 0.30), Inches(1.55)
    )
    add_paragraphs(
        fbox.text_frame,
        [(s, {"size": 11, "space_after": 4, "color": COLOR_TEXT})
         for s in flow_items],
    )

    # ---------- Right wide column: 6 output modules in 2×3 grid ----------
    right_left = 3.85
    right_w = 9.00
    add_filled_rect(slide, right_left, 1.15, right_w, 0.40,
                    fill=COLOR_PRIMARY, corner=0.06)
    add_textbox(
        slide, left=right_left, top=1.15, width=right_w, height=0.40,
        text="6 大模块结构化输出",
        size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    outputs = [
        ("①", "教学目标",
         "可量化、可达成的学习成果"),
        ("②", "重点难点",
         "核心知识点 + 易错环节标注"),
        ("③", "课堂流程",
         "导入 → 讲解 → 演示 → 实训 → 总结"),
        ("④", "基础实训任务",
         "任务驱动 + 步骤指导 + 验收标准"),
        ("⑤", "考核建议",
         "提问 + 实训提交 + 简答评分"),
        ("⑥", "知识来源",
         "知识库片段引用 [1][2][3] 含页码"),
    ]
    grid_top = 1.65
    grid_h = 3.65
    cols = 3
    rows = 2
    cell_gap_x = 0.15
    cell_gap_y = 0.15
    cell_w = (right_w - (cols - 1) * cell_gap_x) / cols
    cell_h = (grid_h - (rows - 1) * cell_gap_y) / rows
    for i, (no, title, desc) in enumerate(outputs):
        col = i % cols
        row = i // cols
        x = right_left + col * (cell_w + cell_gap_x)
        y = grid_top + row * (cell_h + cell_gap_y)
        add_filled_rect(slide, x, y, cell_w, cell_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        # Number badge top-left
        add_filled_rect(slide, x, y, 0.50, cell_h, fill=COLOR_PRIMARY,
                        corner=0.06)
        add_textbox(slide, left=x, top=y, width=0.50, height=cell_h,
                    text=no, size=22, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=x + 0.60, top=y + 0.10,
                    width=cell_w - 0.70, height=0.45,
                    text=title, size=14, bold=True, color=COLOR_PRIMARY)
        add_textbox(slide, left=x + 0.60, top=y + 0.58,
                    width=cell_w - 0.70, height=cell_h - 0.65,
                    text=desc, size=11, color=COLOR_TEXT,
                    line_spacing=1.30)

    # ---------- Reserved screenshot strip notice (5.40 – 5.55) ----------
    add_textbox(
        slide, left=0.55, top=5.42, width=12.30, height=0.22,
        text="◆ 系统演示截图（下方）",
        size=11, bold=True, color=COLOR_ACCENT,
    )

    # Bottom value banner
    add_textbox(
        slide, left=0.55, top=7.10, width=12.30, height=0.32,
        text='备课时间：从"几小时整理素材" → "10 秒生成 + 微调"，输出可直接用于实训课堂',
        size=13, bold=True, color=RGBColor(0xD9, 0x82, 0x2E),
        align=PP_ALIGN.CENTER,
    )


# =============================================================================
# 3. ER picture migration: slide 5 (mistakenly placed) → slide 6 (correct)
# =============================================================================


def migrate_er_picture(prs) -> bool:
    """If a large user picture sits in the lower-left of slide 5 (the
    architecture page), it's almost certainly the ER diagram the user
    intended for slide 6 (the database page). Move it.

    Returns True if a migration happened.
    """
    slides = list(prs.slides)
    if len(slides) < 6:
        return False
    slide5 = slides[4]
    slide6 = slides[5]

    er_pic_shape = None
    for shp in slide5.shapes:
        if shp.shape_type != 13 or is_decoration_logo(shp):
            continue
        w = _in(shp.width); h = _in(shp.height)
        top = _in(shp.top); left = _in(shp.left)
        # Heuristic: large picture in the lower-half left side
        if w >= 4.0 and h >= 2.0 and top >= 3.0 and left <= 6.0:
            er_pic_shape = shp
            break

    if er_pic_shape is None:
        return False

    # Extract image bytes
    blob = er_pic_shape.image.blob

    # Compute aspect-correct fit inside slide 6's ER container.
    # The ER card on slide 6 is left=0.55, top=3.50+0.36=3.86,
    # width=5.40, height=3.40-0.36=3.04 (after the title bar).
    cont_left = 0.60
    cont_top = 3.92
    cont_w = 5.30
    cont_h = 2.95
    try:
        from PIL import Image
        with Image.open(io.BytesIO(blob)) as im:
            iw, ih = im.size
        ratio = iw / ih
    except Exception:
        ratio = 1.0
    cont_ratio = cont_w / cont_h
    if ratio >= cont_ratio:
        draw_w = cont_w
        draw_h = draw_w / ratio
    else:
        draw_h = cont_h
        draw_w = draw_h * ratio
    draw_left = cont_left + (cont_w - draw_w) / 2
    draw_top = cont_top + (cont_h - draw_h) / 2

    slide6.shapes.add_picture(
        io.BytesIO(blob),
        Inches(draw_left), Inches(draw_top),
        width=Inches(draw_w), height=Inches(draw_h),
    )

    # Remove from slide 5
    sp = er_pic_shape._element
    sp.getparent().remove(sp)
    return True


def rebuild_remaining_work(slide) -> None:
    """Remaining-work page (8 items, with industrial-grade highlight item ④).

    Layout: 8 work items on the left + 8-row W9~W16 Gantt on the right
    + bottom banner + a small thesis-mapping caption beneath the banner.

    No risk-mitigation block (intentionally removed)."""
    clear_slide_keep_logo_only(slide)
    add_title_bar(slide, "剩余工作计划（第 9 周 → 第 16 周）")

    # Subtitle
    add_textbox(
        slide, left=0.55, top=1.05, width=12.30, height=0.30,
        text="◆ 后续 8 周工作安排",
        size=14, bold=True, color=COLOR_PRIMARY,
    )

    # 8 work items (item ④ is the industrial-grade highlight)
    items = [
        ("①", "主观题高级评分",
         "LLM-as-Judge（Qwen-Max 评委）+ 语义相似度加权融合"),
        ("②", "学情分析",
         "错题聚合、班级统计、薄弱知识点识别 + ECharts"),
        ("③", "管理端数据看板",
         "多维图表：参与人数 / 正确率 / 知识点掌握分布"),
        ("④", "评测体系 + 可观测性 ★",
         "Golden Set 100 QA + RAGAS 4 指标 + LangFuse Trace"),
        ("⑤", "系统稳定性 + 性能",
         "异常处理、SSE 重连、Redis 缓存加速、压力测试"),
        ("⑥", "系统测试",
         "功能 + 性能 + 用户测试 + CI 自动回归"),
        ("⑦", "论文撰写与查重",
         "毕业设计论文（含评测实验章节）+ 学校查重"),
        ("⑧", "PPT + 预答辩 + 修改",
         "正式答辩材料整理"),
    ]
    base_top = 1.50
    item_h = 0.55
    gap_y = 0.05
    item_w = 5.30
    item_left = 0.55
    HIGHLIGHT_COLOR = RGBColor(0xD9, 0x82, 0x2E)  # orange for the spotlight item ④
    for i, (no, title, desc) in enumerate(items):
        y = base_top + i * (item_h + gap_y)
        is_highlight = (i == 3)
        bar_color = HIGHLIGHT_COLOR if is_highlight else COLOR_PRIMARY
        bg_color = (RGBColor(0xFD, 0xF1, 0xE2) if is_highlight else COLOR_LIGHT_BG)
        add_filled_rect(slide, item_left, y, item_w, item_h,
                        fill=bg_color, corner=0.04)
        add_filled_rect(slide, item_left, y, 0.42, item_h,
                        fill=bar_color, corner=0.04)
        add_textbox(slide, left=item_left, top=y, width=0.42, height=item_h,
                    text=no, size=14, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=item_left + 0.50, top=y + 0.02,
                    width=1.85, height=item_h - 0.04,
                    text=title, size=11, bold=True,
                    color=bar_color,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left=item_left + 2.40, top=y + 0.02,
                    width=item_w - 2.50, height=item_h - 0.04,
                    text=desc, size=9, color=COLOR_TEXT,
                    anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20)

    # === Gantt chart (right) ===
    g_left = 6.20
    g_w = 6.65
    g_top = 1.50
    add_filled_rect(slide, g_left, g_top, g_w, 0.36, fill=COLOR_PRIMARY,
                    corner=0.04)
    add_textbox(
        slide, left=g_left, top=g_top, width=g_w, height=0.36,
        text="◆ 时间安排（W9 ~ W16）", size=12, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Week column headers
    week_top = g_top + 0.42
    weeks = ["W9", "W10", "W11", "W12", "W13", "W14", "W15", "W16"]
    n_weeks = len(weeks)
    week_w = g_w / n_weeks
    for i, w in enumerate(weeks):
        x = g_left + i * week_w
        add_filled_rect(slide, x, week_top, week_w - 0.02, 0.26,
                        fill=COLOR_GRAY_BG, corner=0.02)
        add_textbox(slide, left=x, top=week_top, width=week_w - 0.02,
                    height=0.26, text=w, size=9, bold=True,
                    color=COLOR_PRIMARY, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, font_name=EN_FONT)

    # Bars: (item_idx_zero_based, start_week, end_week, color)
    bars = [
        (0, 0, 1, RGBColor(0x4A, 0x90, 0xE2)),  # ① 主观题评分    W9-W10
        (1, 1, 2, RGBColor(0x2E, 0x9C, 0x6E)),  # ② 学情分析      W10-W11
        (2, 1, 2, RGBColor(0xE5, 0xA5, 0x3B)),  # ③ 数据看板      W10-W11
        (3, 2, 3, HIGHLIGHT_COLOR),              # ④ 评测体系 ★    W11-W12
        (4, 3, 4, RGBColor(0xC0, 0x55, 0x55)),  # ⑤ 稳定性+性能   W12-W13
        (5, 4, 5, RGBColor(0x69, 0x40, 0x90)),  # ⑥ 系统测试      W13-W14
        (6, 4, 6, RGBColor(0x2E, 0x6F, 0xB5)),  # ⑦ 论文撰写      W13-W15
        (7, 6, 7, RGBColor(0x1E, 0x40, 0x80)),  # ⑧ PPT + 预答辩  W15-W16
    ]
    bar_top0 = week_top + 0.32
    bar_h = 0.42  # 8 bars × (0.42+0.04) = 3.68 → ends at 4.74
    bar_gap = 0.04
    no_glyphs = ['①', '②', '③', '④', '⑤', '⑥', '⑦', '⑧']
    for idx, sw, ew, color in bars:
        y = bar_top0 + idx * (bar_h + bar_gap)
        x = g_left + sw * week_w + 0.04
        w = (ew - sw + 1) * week_w - 0.08
        add_filled_rect(slide, x, y, w, bar_h, fill=color, corner=0.06)
        add_textbox(slide, left=x + 0.06, top=y, width=w - 0.12,
                    height=bar_h,
                    text=f"{no_glyphs[idx]}  {items[idx][1]}",
                    size=9, bold=True,
                    color=RGBColor(0xFF, 0xFF, 0xFF),
                    anchor=MSO_ANCHOR.MIDDLE)

    # ===== Bottom: banner + thesis-mapping note =====
    banner_top = 6.95
    add_filled_rect(slide, 0.55, banner_top, 12.30, 0.30,
                    fill=COLOR_PRIMARY,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                    corner=0.40)
    add_textbox(
        slide, left=0.55, top=banner_top, width=12.30, height=0.30,
        text="任务书剩余项 + 论文按周推进；评测 / 可观测性作为论文核心实验章节，提供量化数据支撑",
        size=12, bold=True, color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Thesis-mapping note (small caption beneath the banner)
    add_textbox(
        slide, left=0.55, top=banner_top + 0.32, width=12.30, height=0.18,
        text="★ 第 ④ 项「评测体系 + 可观测性」对应论文「系统评测」章节，提供 Golden Set 上的量化对比数据，"
             "让\"系统能用\"升级为\"系统好用、可衡量\"",
        size=9, color=RGBColor(0x66, 0x66, 0x66),
        align=PP_ALIGN.CENTER,
    )


# =============================================================================
# Additive helpers (purely add new shapes, do NOT clear / mutate existing ones)
# =============================================================================


def add_p4_scope_note(slide) -> None:
    """Append a one-line scope note to slide 4 (architecture).

    Purely additive — DOES NOT clear or mutate any existing shape on the
    slide. The 8-tier architecture and the dual closed-loop are preserved
    as-is. The note sits in the small gap between the 8-tier main area
    (ends near 6.56") and the dual closed-loop (starts near 6.78").
    """
    add_textbox(
        slide, left=0.55, top=6.58, width=12.30, height=0.16,
        text="★ 中期已落地：FastAPI 业务层 + 数据层 + LangChain 编排 + DashScope AI 模型层 + 三角色 RBAC ；"
             "阶段四（W10-W13）渐进推进：Nginx 网关 · LangGraph + MCP + Skills · Redis 中间件 · Ollama 兜底 · "
             "SpringBoot 微服务 · LangFuse + Golden Set + CI Gate",
        size=8, color=RGBColor(0x66, 0x66, 0x66),
        align=PP_ALIGN.CENTER,
    )


def duplicate_slide_after(prs, source_idx: int):
    """Duplicate the slide at `source_idx` and insert the copy at
    position `source_idx + 1`.

    Returns the new slide object. Implemented via XML-level deep copy
    of all shapes + re-relating image/oleObject relationships, since
    python-pptx does not provide a built-in slide-copy API.
    """
    from copy import deepcopy

    source = prs.slides[source_idx]
    blank_layout = source.slide_layout
    new_slide = prs.slides.add_slide(blank_layout)

    # Copy each shape from source via XML deep-copy
    src_tree = source.shapes._element
    new_tree = new_slide.shapes._element
    for shp in source.shapes:
        new_tree.append(deepcopy(shp.element))

    # Re-establish relationships (images, etc.) so picture rId's resolve
    for rel in source.part.rels.values():
        if "notesSlide" in rel.reltype:
            continue
        try:
            new_slide.part.relate_to(rel.target_part, rel.reltype)
        except Exception:
            # If re-relating fails for an unusual rel type, skip silently
            pass

    # Move the new slide from end-of-deck to source_idx + 1
    sld_id_lst = prs.slides._sldIdLst
    sld_ids = list(sld_id_lst)
    new_id = sld_ids[-1]
    sld_id_lst.remove(new_id)
    sld_id_lst.insert(source_idx + 1, new_id)

    return new_slide


def rebuild_p6_login(slide) -> None:
    """Build the new "用户登录与角色管理" page (slide 6 after insertion).

    Layout (aligned to 中期答辩PPT大纲.md):
      - Title bar
      - Top 承接说明 banner (one-line, light)
      - Middle: left flow (registration + login) | right RBAC matrix
      - 下方: 4 security highlight cards in a single row
      - 底部: 演示截图占位符 + 金句横幅
    """
    clear_slide_keep_logo_only(slide)
    add_title_bar(slide, "用户登录与角色管理")

    # ---------------- Top: 承接说明 ----------------
    add_filled_rect(slide, 0.55, 1.05, 12.30, 0.36, fill=COLOR_LIGHT_BG, corner=0.04)
    add_textbox(
        slide, left=0.55, top=1.05, width=12.30, height=0.36,
        text="◆ 承接第 5 页 — 基于 users + sessions 表实现统一认证体系，配合三角色 RBAC 路由守卫",
        size=12, bold=True, color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ===================================================================
    # Middle row : left flow (5.50") | right RBAC matrix (6.65")
    # Compressed from 4.0" → 3.00" to make room for the screenshot strip.
    # ===================================================================
    middle_top = 1.55
    middle_h = 3.00

    # ---------------- Left: 登录注册流程 ----------------
    flow_left = 0.55
    flow_w = 5.50
    add_filled_rect(slide, flow_left, middle_top, flow_w, 0.36,
                    fill=COLOR_PRIMARY, corner=0.04)
    add_textbox(
        slide, left=flow_left, top=middle_top, width=flow_w, height=0.36,
        text="◆ 登录 / 注册流程", size=12, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # 5 stacked steps (vertical pipeline)
    steps = [
        ("注册", "邮箱 + 密码 + 角色 (teacher / student / admin)", COLOR_ACCENT),
        ("哈希", "PBKDF2-HMAC-SHA256 · 10 万次迭代 · 16 字节随机 salt",
         RGBColor(0xE5, 0xA5, 0x3B)),
        ("入库", "INSERT users (password_hash / password_salt 分字段存储)",
         RGBColor(0x2E, 0x9C, 0x6E)),
        ("登录", "查 users → 同 salt 重新 hash → 比对 → 写 sessions",
         COLOR_ACCENT),
        ("鉴权", "前端 localStorage → Authorization Bearer → require_user 注入",
         RGBColor(0xC0, 0x55, 0x55)),
    ]
    n_steps = len(steps)
    flow_body_top = middle_top + 0.46
    flow_body_h = middle_h - 0.50
    step_h = (flow_body_h - (n_steps - 1) * 0.08) / n_steps
    for i, (head, desc, color) in enumerate(steps):
        y = flow_body_top + i * (step_h + 0.08)
        # Step pill
        add_filled_rect(slide, flow_left, y, flow_w, step_h,
                        fill=COLOR_LIGHT_BG, corner=0.05)
        add_filled_rect(slide, flow_left, y, 1.10, step_h, fill=color, corner=0.05)
        add_textbox(
            slide, left=flow_left, top=y, width=1.10, height=step_h,
            text=head, size=12, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left=flow_left + 1.20, top=y + 0.04,
            width=flow_w - 1.30, height=step_h - 0.08,
            text=desc, size=10, color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.20,
        )
        # Down arrow between steps
        if i < n_steps - 1:
            arrow_y = y + step_h + 0.005
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(flow_left + 0.50), Inches(arrow_y),
                Inches(0.18), Inches(0.07),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_SUBTLE
            arrow.line.fill.background()

    # ---------------- Right: RBAC matrix ----------------
    rbac_left = 6.20
    rbac_w = 6.65
    add_filled_rect(slide, rbac_left, middle_top, rbac_w, 0.36,
                    fill=RGBColor(0xD9, 0x82, 0x2E), corner=0.04)
    add_textbox(
        slide, left=rbac_left, top=middle_top, width=rbac_w, height=0.36,
        text="◆ 三角色 RBAC 权限矩阵", size=12, bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    roles = [
        {
            "name": "教师",
            "en": "teacher",
            "color": RGBColor(0x4A, 0x90, 0xE2),
            "routes": "/courses · /courses/new · /knowledge-base · "
                      "/lesson-outline · /exercises",
            "ability": "知识库管理 / 课程创建 / 备课 / 出题 / 学情分析",
        },
        {
            "name": "学生",
            "en": "student",
            "color": RGBColor(0x2E, 0x9C, 0x6E),
            "routes": "/courses · /qa · /exercises/session · /personalized",
            "ability": "RAG 问答 / 完成练习 / 查看掌握度 / 个性化推荐",
        },
        {
            "name": "管理员",
            "en": "admin",
            "color": RGBColor(0xC0, 0x55, 0x55),
            "routes": "/admin · /knowledge-base · /courses",
            "ability": "用户管理 / 课程管理 / 数据看板 / 全局知识库",
        },
    ]
    rb_body_top = middle_top + 0.46
    rb_body_h = middle_h - 0.50
    role_h = (rb_body_h - (len(roles) - 1) * 0.10) / len(roles)
    for i, r in enumerate(roles):
        y = rb_body_top + i * (role_h + 0.10)
        # Card
        add_filled_rect(slide, rbac_left, y, rbac_w, role_h,
                        fill=COLOR_LIGHT_BG, corner=0.05)
        # Role label band (left)
        add_filled_rect(slide, rbac_left, y, 1.20, role_h,
                        fill=r["color"], corner=0.05)
        add_textbox(
            slide, left=rbac_left, top=y + 0.10, width=1.20, height=role_h * 0.40,
            text=r["name"], size=14, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left=rbac_left, top=y + role_h * 0.50,
            width=1.20, height=role_h * 0.35,
            text=r["en"], size=9,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, font_name=EN_FONT,
        )
        # Routes & ability (right)
        add_textbox(
            slide, left=rbac_left + 1.30, top=y + 0.06,
            width=rbac_w - 1.40, height=role_h * 0.45,
            text=f"路由：{r['routes']}", size=10, color=COLOR_PRIMARY,
            line_spacing=1.20,
        )
        add_textbox(
            slide, left=rbac_left + 1.30, top=y + role_h * 0.50,
            width=rbac_w - 1.40, height=role_h * 0.45,
            text=f"能力：{r['ability']}", size=10, color=COLOR_TEXT,
            line_spacing=1.20,
        )

    # ===================================================================
    # 下方: 4 security highlight cards (single row, compact)
    # ===================================================================
    sec_top = middle_top + middle_h + 0.20  # ≈ 4.75
    sec_h = 0.78
    add_textbox(
        slide, left=0.55, top=sec_top - 0.26, width=12.30, height=0.22,
        text="◆ 安全设计要点", size=11, bold=True, color=COLOR_PRIMARY,
    )
    sec_items = [
        ("PBKDF2-HMAC-SHA256",
         "NIST 推荐密码哈希 · 10 万次迭代抗暴力破解"),
        ("独立 Salt",
         "每用户 16 字节随机 salt · 彻底防彩虹表攻击"),
        ("hash / salt 分存",
         "password_hash 与 password_salt 分字段存储不可还原"),
        ("双层路由守卫",
         "前端 router.beforeEach (UX) · 后端 require_user 依赖注入"),
    ]
    n_sec = len(sec_items)
    sec_gap = 0.15
    sec_w = (12.30 - (n_sec - 1) * sec_gap) / n_sec  # ≈ 2.96
    for i, (head, desc) in enumerate(sec_items):
        x = 0.55 + i * (sec_w + sec_gap)
        add_filled_rect(slide, x, sec_top, sec_w, sec_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, x, sec_top, sec_w, 0.28,
                        fill=COLOR_PRIMARY, corner=0.06)
        add_textbox(
            slide, left=x, top=sec_top + 0.01, width=sec_w, height=0.26,
            text=head, size=11, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left=x + 0.10, top=sec_top + 0.30, width=sec_w - 0.20,
            height=sec_h - 0.34, text=desc, size=9, color=COLOR_TEXT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.20,
        )

    # ===================================================================
    # 底部: 演示截图占位符 (single row, full width)
    # ===================================================================
    shot_top = sec_top + sec_h + 0.18  # ≈ 5.71
    shot_h = 1.30
    shot_left = 0.55
    shot_w = 12.30
    add_filled_rect(slide, shot_left, shot_top, shot_w, 0.28,
                    fill=COLOR_ACCENT, corner=0.04)
    add_textbox(
        slide, left=shot_left, top=shot_top, width=shot_w, height=0.28,
        text="◆ 演示截图  ·  /login  Login.vue",
        size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    add_image_placeholder(
        slide,
        left=shot_left,
        top=shot_top + 0.30,
        width=shot_w,
        height=shot_h - 0.30,
        label='登录 / 注册页（左右分栏品牌设计） — Tab 切换 "账户登录 / 新用户注册"\n'
              'feature badges：RAG 增强问答 · 自动练习生成 · 实时学情分析',
    )

    # ---------------- 底部 golden banner ----------------
    banner_top = shot_top + shot_h + 0.05  # ≈ 7.06
    add_filled_rect(slide, 0.55, banner_top, 12.30, 0.30,
                    fill=COLOR_PRIMARY,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                    corner=0.40)
    add_textbox(
        slide, left=0.55, top=banner_top, width=12.30, height=0.30,
        text="一套统一的 PBKDF2 + 三角色 RBAC 认证体系，支撑教师 / 学生 / 管理三端的安全访问",
        size=12, bold=True, color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# =============================================================================
# Main
# =============================================================================


def main() -> int:
    """Plan-aligned incremental edit. Only the THREE plan-tagged spots
    are touched; every other slide (incl. user-pasted screenshots and
    the ER diagram) is left exactly as-is. The global font_size +1pt
    and emoji cleanup passes are NOT re-run (they would re-grow fonts).
    """
    if not PPT_PATH.exists():
        print(f"[ERR] PPT not found: {PPT_PATH}")
        return 1

    prs = Presentation(PPT_PATH)
    n_before = len(prs.slides)
    print(f"Loaded: {PPT_PATH.name}  ({n_before} slides)")

    # ------------------------------------------------------------------
    # Step 1 — Insert a new blank slide right after slide 5 (database).
    # If the deck is already 13 pages the insertion is skipped (idempotent).
    # ------------------------------------------------------------------
    if n_before == 12:
        new_slide = duplicate_slide_after(prs, source_idx=4)  # after slide 5
        print("  [insert ] duplicated slide 5 → new blank container at slide 6")
    elif n_before == 13:
        new_slide = list(prs.slides)[5]
        print("  [skip   ] deck already 13 slides; reusing existing slide 6 as login page")
    else:
        print(f"  [WARN   ] unexpected deck size {n_before}; aborting login insertion")
        new_slide = None

    # ------------------------------------------------------------------
    # Step 2 — Build the new "用户登录与角色管理" page on slide 6.
    # ------------------------------------------------------------------
    if new_slide is not None:
        rebuild_p6_login(new_slide)
        print("  [rebuilt] slide 6 (login)  fresh build")

    # ------------------------------------------------------------------
    # Step 3 — Add a small scope note to slide 4 (architecture).
    # Purely additive — does NOT modify any existing shape.
    # ------------------------------------------------------------------
    slides = list(prs.slides)
    if len(slides) >= 4:
        add_p4_scope_note(slides[3])
        print("  [appended] slide 4 (architecture)  scope note added at top:6.58\"")

    # ------------------------------------------------------------------
    # Step 4 — Rewrite the remaining-work page (now slide 12 after the
    # insertion of slide 6 above). Falls back to slide 11 if no insertion
    # was performed (deck stayed at 12).
    # ------------------------------------------------------------------
    remaining_idx = 11 if len(slides) >= 13 else 10
    if remaining_idx < len(slides):
        before_pics = sum(1 for s in slides[remaining_idx].shapes
                          if is_user_picture(s))
        rebuild_remaining_work(slides[remaining_idx])
        after_pics = sum(1 for s in slides[remaining_idx].shapes
                         if is_user_picture(s))
        print(f"  [rebuilt] slide {remaining_idx + 1} (remaining-work)  "
              f"8 items + 论文对位 (pics {before_pics} → {after_pics})")

    prs.save(PPT_PATH)
    print(f"Saved: {PPT_PATH.name}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
