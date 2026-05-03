"""One-off: re-build ONLY slide 6 (用户登录与角色管理) to align with
中期答辩PPT大纲.md (adds 演示截图占位符 in the 底部 row).

Does NOT touch any other slide — explicitly avoids the global passes
and the additive helpers from `update_defense_pptx.main`.

Run:
    .ppt_venv/bin/python scripts/update_p6_only.py
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from update_defense_pptx import (  # noqa: E402
    PPT_PATH,
    is_user_picture,
    rebuild_p6_login,
)
from pptx import Presentation  # noqa: E402


def main() -> int:
    if not PPT_PATH.exists():
        print(f"[ERR] PPT not found: {PPT_PATH}")
        return 1

    prs = Presentation(PPT_PATH)
    slides = list(prs.slides)
    print(f"Loaded: {PPT_PATH.name}  ({len(slides)} slides)")

    if len(slides) < 6:
        print(f"[ERR] deck has only {len(slides)} slides; expected ≥ 6")
        return 2

    slide6 = slides[5]
    user_pics_before = sum(1 for s in slide6.shapes if is_user_picture(s))
    rebuild_p6_login(slide6)
    user_pics_after = sum(1 for s in slide6.shapes if is_user_picture(s))
    print(
        f"  [rebuilt] slide 6 (login)  fresh build  "
        f"(user_pics {user_pics_before} → {user_pics_after})"
    )

    prs.save(PPT_PATH)
    print(f"Saved: {PPT_PATH.name}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
