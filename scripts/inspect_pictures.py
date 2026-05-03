"""Inspect picture sizes in template slides to identify decorations vs content."""
from pathlib import Path
from pptx import Presentation

TEMPLATE_PATH = Path(
    "/Users/daniel/答辩PPT/中国石油大学（华东）PPT模板/学海遨游，逐梦石大-学术系列_副本.pptx"
)


def main() -> None:
    prs = Presentation(TEMPLATE_PATH)
    for si, slide in enumerate(prs.slides):
        if si > 13 and si < 26:
            continue
        print(f"\n=== Slide {si + 1} ({slide.slide_layout.name}) ===")
        for shp in slide.shapes:
            if shp.shape_type != 13:
                continue
            try:
                w = shp.width / 914400
                h = shp.height / 914400
                top = shp.top / 914400
                left = shp.left / 914400
                area = (w * h) / (13.33 * 7.5)
                print(
                    f"  PIC name={shp.name!r:20s} pos=({left:5.2f},{top:5.2f}) "
                    f"size=({w:5.2f} x {h:5.2f}) area={area:.2f}"
                )
            except Exception as e:
                print(f"  PIC name={shp.name!r} ERROR={e}")


if __name__ == "__main__":
    main()
