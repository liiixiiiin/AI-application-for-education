"""
Build the mid-term defense PPT based on the UPC academic template.

Strategy:
- Load the template (preserves the blue-curve background, school logo,
  and other decorations which are baked into each individual slide).
- Reuse the template's existing slides as canvases:
    Slide 1  -> our cover page (keep all decorations, clear text)
    Slides 2-13 -> our 12 content pages (strip text/content shapes,
                   redraw our content on top, keep PICTURE backgrounds)
    Slide 27 -> our thanks page (keep decorations, replace text)
- Delete the unused template slides (14-26).

Output: /Users/daniel/2026design/App/中期答辩PPT.pptx
"""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt, Emu, Inches

# ------------------------------------------------------------------
# Paths & constants
# ------------------------------------------------------------------
TEMPLATE_PATH = Path(
    "/Users/daniel/答辩PPT/中国石油大学（华东）PPT模板/学海遨游，逐梦石大-学术系列_副本.pptx"
)
OUTPUT_PATH = Path("/Users/daniel/2026design/App/中期答辩PPT.pptx")
RESOURCES = Path("/Users/daniel/2026design/App/ppt_resources")
ER_DIAGRAM_PNG = Path("/Users/daniel/2026design/App/data/er-diagram/er.png")

# UPC blue color palette (eyeballed from the template preview)
COLOR_PRIMARY = RGBColor(0x1E, 0x40, 0x80)  # deep blue
COLOR_ACCENT = RGBColor(0x2E, 0x6F, 0xB5)    # mid blue
COLOR_LIGHT_BG = RGBColor(0xEA, 0xF2, 0xFB)  # very light blue
COLOR_GRAY_BG = RGBColor(0xF4, 0xF6, 0xF9)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)
COLOR_SUBTLE = RGBColor(0x88, 0x88, 0x88)
COLOR_PLACEHOLDER_BORDER = RGBColor(0x9F, 0xB6, 0xD4)

CN_FONT = "微软雅黑"
EN_FONT = "Calibri"

# Slide dimensions: 13.333" x 7.5" (16:9 widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ------------------------------------------------------------------
# Helpers
# ------------------------------------------------------------------
def delete_slide_at(prs: Presentation, slide_index: int) -> None:
    """Delete a single slide at the given index from the presentation."""
    sldIdLst = prs.slides._sldIdLst  # noqa: SLF001
    sldIds = list(sldIdLst)
    if slide_index >= len(sldIds):
        return
    sldId = sldIds[slide_index]
    rid = sldId.get(qn("r:id"))
    sldIdLst.remove(sldId)
    rels = prs.part.rels
    try:
        rels._rels.pop(rid, None)  # noqa: SLF001
    except Exception:
        pass


def clear_slide_content(slide, *, keep_pictures: bool = True) -> None:
    """Strip text boxes and content rectangles from a slide while keeping
    decorative pictures (background, logo) and connectors.

    - Removes: TEXT_BOX (17), AUTO_SHAPE (1), FREEFORM (5) shapes that aren't
      decorations, GROUP (6) (template grouping that includes textual labels).
    - Keeps: PICTURE (13) decorations and LINE (9) connectors.
    """
    spTree = slide.shapes._spTree  # noqa: SLF001
    shapes_to_remove = []
    for shp in slide.shapes:
        st = shp.shape_type
        # Always remove text boxes
        if st == 17:  # TEXT_BOX
            shapes_to_remove.append(shp)
            continue
        # Remove rectangles / freeforms (likely content cards or content shapes)
        if st in (1, 5, 6):  # AUTO_SHAPE, FREEFORM, GROUP
            shapes_to_remove.append(shp)
            continue
        # Remove placeholders that have text
        if st == 14:  # PLACEHOLDER
            shapes_to_remove.append(shp)
            continue
        # Optionally remove pictures (e.g., scenery on content pages)
        if not keep_pictures and st == 13:
            shapes_to_remove.append(shp)
    for shp in shapes_to_remove:
        sp = shp._element  # noqa: SLF001
        sp.getparent().remove(sp)


def clear_slide_for_content_page(slide) -> None:
    """For content pages: remove text/content shapes and large content
    pictures (scenery, campus photos), keep only the small school logo at
    the top-right. The template's grid background is a slide background
    fill (not a picture) so it's preserved automatically.
    """
    # First remove text/content shapes
    clear_slide_content(slide, keep_pictures=True)
    # Remove all pictures except the small logo in upper-right
    pictures_to_remove = []
    for shp in slide.shapes:
        if shp.shape_type != 13:  # not picture
            continue
        try:
            w = shp.width / 914400
            top = shp.top / 914400
            left = shp.left / 914400
        except Exception:
            continue
        # Keep only the small logo (top < 1.0, top-right corner, width < 3)
        if top < 1.0 and left > 9.5 and w < 3.0:
            continue
        pictures_to_remove.append(shp)
    for shp in pictures_to_remove:
        sp = shp._element  # noqa: SLF001
        sp.getparent().remove(sp)


def set_text(
    text_frame,
    text: str,
    *,
    font_name: str = CN_FONT,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float | None = None,
) -> None:
    """Replace the contents of a text frame with a single styled paragraph."""
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor
    p = text_frame.paragraphs[0]
    p.alignment = align
    if line_spacing is not None:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_paragraphs(
    text_frame,
    items: list[tuple[str, dict]],
    *,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> None:
    """Add multiple styled paragraphs to a text frame.

    Each item is (text, style_dict). style_dict supports: size, bold, color,
    font_name, space_before, space_after, line_spacing, indent.
    """
    text_frame.clear()
    text_frame.word_wrap = True
    text_frame.vertical_anchor = anchor

    for i, (text, style) in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.alignment = style.get("align", align)
        if "space_before" in style:
            p.space_before = Pt(style["space_before"])
        if "space_after" in style:
            p.space_after = Pt(style["space_after"])
        if "line_spacing" in style:
            p.line_spacing = style["line_spacing"]
        if "level" in style:
            p.level = style["level"]
        run = p.add_run()
        run.text = text
        run.font.name = style.get("font_name", CN_FONT)
        run.font.size = Pt(style.get("size", 16))
        run.font.bold = style.get("bold", False)
        run.font.color.rgb = style.get("color", COLOR_TEXT)


def add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str = "",
    *,
    size: int = 16,
    bold: bool = False,
    color: RGBColor = COLOR_TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    font_name: str = CN_FONT,
    line_spacing: float | None = None,
):
    """Add a textbox at given inch coordinates."""
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    if text:
        set_text(
            tf,
            text,
            size=size,
            bold=bold,
            color=color,
            align=align,
            anchor=anchor,
            font_name=font_name,
            line_spacing=line_spacing,
        )
    else:
        tf.clear()
    return box


def add_filled_rect(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    fill: RGBColor = COLOR_LIGHT_BG,
    line: RGBColor | None = None,
    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
    corner: float = 0.05,
):
    """Add a filled rounded-rectangle background shape."""
    rect = slide.shapes.add_shape(
        shape_type, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(0.75)
    rect.shadow.inherit = False
    # Adjust corner radius if applicable
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            rect.adjustments[0] = corner
        except Exception:
            pass
    return rect


def add_image_placeholder(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    label: str = "在此处插入图片",
):
    """Draw a dashed rectangle placeholder with hint text for an image."""
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_GRAY_BG
    rect.line.color.rgb = COLOR_PLACEHOLDER_BORDER
    rect.line.width = Pt(1.5)
    # set dashed line via lxml
    ln = rect.line._get_or_add_ln()  # noqa: SLF001
    prstDash = ln.find(qn("a:prstDash"))
    if prstDash is None:
        from lxml import etree

        prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")

    tf = rect.text_frame
    tf.word_wrap = True
    set_text(
        tf,
        label,
        size=14,
        color=COLOR_SUBTLE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    return rect


def add_title_bar(slide, title: str, number: str | None = None) -> None:
    """Add the section title bar at the top-left of every content slide.

    Mimics the template style:
      [N、 标题文字] with a deep-blue color. School logo & decorations are
    inherited from the master so we don't draw them here.
    """
    # Vertical accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(0.45), Inches(0.10), Inches(0.40)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLOR_PRIMARY
    accent.line.fill.background()

    # Title text
    text = f"{number}、{title}" if number else title
    add_textbox(
        slide,
        left=0.65,
        top=0.40,
        width=10.5,
        height=0.55,
        text=text,
        size=26,
        bold=True,
        color=COLOR_PRIMARY,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Thin underline below the title
    line = slide.shapes.add_connector(
        1, Inches(0.45), Inches(1.05), Inches(12.85), Inches(1.05)
    )
    line.line.color.rgb = COLOR_PRIMARY
    line.line.width = Pt(1.2)


def add_page_number(slide, idx: int, total: int = 13) -> None:
    add_textbox(
        slide,
        left=12.0,
        top=7.10,
        width=1.2,
        height=0.30,
        text=f"{idx} / {total}",
        size=10,
        color=COLOR_SUBTLE,
        align=PP_ALIGN.RIGHT,
    )


# ------------------------------------------------------------------
# Slide builders
# ------------------------------------------------------------------
def build_cover(slide) -> None:
    """Slide 1: Cover. Operates on the template's existing cover slide,
    keeping the blue-curve decoration & school logo, removing the
    "学海遨游 逐梦石大" artistic title image and other text shapes."""
    to_remove = []
    for shp in slide.shapes:
        st = shp.shape_type
        # Remove text boxes, placeholders, and group shapes
        if st in (17, 14, 6):  # TEXT_BOX, PLACEHOLDER, GROUP
            to_remove.append(shp)
            continue
        if st == 1:  # AUTO_SHAPE - check if it has text
            try:
                if shp.has_text_frame and shp.text_frame.text.strip():
                    to_remove.append(shp)
                    continue
            except Exception:
                pass
        # Remove the wide artistic title image (area > 0.10, wide horizontal)
        if st == 13:  # PICTURE
            try:
                w = shp.width / 914400
                h = shp.height / 914400
                top = shp.top / 914400
                # Wide horizontal banner image (the artistic title)
                if w > 8.0 and h < 2.0 and top > 1.5 and top < 4.5:
                    to_remove.append(shp)
            except Exception:
                pass
    for shp in to_remove:
        sp = shp._element  # noqa: SLF001
        sp.getparent().remove(sp)

    # Main title
    add_textbox(
        slide,
        left=1.5,
        top=2.30,
        width=10.3,
        height=1.2,
        text="基于 AI 大模型的教学实训智能辅助系统的设计与实现",
        size=36,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Subtitle (English)
    add_textbox(
        slide,
        left=1.5,
        top=3.50,
        width=10.3,
        height=0.5,
        text="Design and Implementation of an AI-LLM-Based Intelligent Teaching Training Assistant System",
        size=14,
        color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
        font_name=EN_FONT,
    )
    # Tag line
    add_textbox(
        slide,
        left=1.5,
        top=4.10,
        width=10.3,
        height=0.5,
        text="—— 本科毕业设计 中期答辩 ——",
        size=18,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )

    # Info block (centered)
    info = [
        ("学    生：李    鑫", {"size": 18, "color": COLOR_TEXT, "space_after": 8}),
        ("学    号：2207020422", {"size": 18, "color": COLOR_TEXT, "space_after": 8}),
        ("指导教师：XXX", {"size": 18, "color": COLOR_TEXT, "space_after": 8}),
        ("学    院：计算机科学与技术学院", {"size": 18, "color": COLOR_TEXT, "space_after": 8}),
        ("答辩日期：2026 年 X 月 X 日", {"size": 18, "color": COLOR_TEXT}),
    ]
    box = slide.shapes.add_textbox(Inches(4.3), Inches(5.10), Inches(4.7), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    add_paragraphs(tf, info, align=PP_ALIGN.LEFT)


def build_section_divider(
    prs: Presentation, number: str, title: str, subtitle: str = ""
) -> None:
    """Optional: section divider. Not used for the 14-page version, but kept
    in case we want to insert section breaks later."""
    layout = prs.slide_masters[0].slide_layouts[0]  # 空白
    slide = prs.slides.add_slide(layout)
    add_filled_rect(
        slide, 0.45, 3.0, 12.4, 1.5, fill=COLOR_PRIMARY, shape_type=MSO_SHAPE.RECTANGLE
    )
    add_textbox(
        slide,
        left=0.45,
        top=3.0,
        width=12.4,
        height=1.5,
        text=f"{number}  {title}",
        size=40,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    if subtitle:
        add_textbox(
            slide,
            left=0.45,
            top=4.55,
            width=12.4,
            height=0.6,
            text=subtitle,
            size=18,
            color=COLOR_PRIMARY,
            align=PP_ALIGN.CENTER,
        )


def prepare_content_slide(slide):
    """Prepare an existing template content slide as our canvas: strip
    text/content shapes, keep background grid + school logo."""
    clear_slide_for_content_page(slide)
    return slide


def build_p2_background(slide) -> None:
    """Page 2: 选题背景.

    Layout: left column shows two policy document images (Ministry of
    Education + State Council) as authoritative backing; right column uses
    a STAR-like narrative (Situation - Task/Pain - Action/Innovation -
    Result/Value) to grab the panel's attention.
    """
    prepare_content_slide(slide)
    add_title_bar(slide, "选题背景", "一")

    # =========================================================
    # Left column: two policy document screenshots (5.0" wide)
    # =========================================================
    LEFT_X = 0.55
    LEFT_W = 4.45

    # Section heading on the left
    add_textbox(
        slide,
        left=LEFT_X,
        top=1.20,
        width=LEFT_W,
        height=0.40,
        text="◆ 政策权威背书",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # Top image: 教育部数字化 (aspect ratio ≈ 1.732)
    img1_top = 1.65
    img1_h = LEFT_W / 1.732  # ≈ 2.57
    img1_h = 2.20  # constrain a bit shorter to leave room for caption
    img1_w = img1_h * 1.732  # ≈ 3.81
    img1_left = LEFT_X + (LEFT_W - img1_w) / 2  # center
    pic1 = slide.shapes.add_picture(
        str(RESOURCES / "教育部数字化.jpg"),
        Inches(img1_left),
        Inches(img1_top),
        width=Inches(img1_w),
        height=Inches(img1_h),
    )
    # Subtle border around the image
    pic1.line.color.rgb = COLOR_PRIMARY
    pic1.line.width = Pt(0.75)

    # Caption 1
    add_textbox(
        slide,
        left=LEFT_X,
        top=img1_top + img1_h + 0.05,
        width=LEFT_W,
        height=0.32,
        text="教育部 ·《教育数字化战略行动》",
        size=11,
        bold=True,
        color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
    )

    # Bottom image: 国务院人工智能行动 (aspect ratio ≈ 2.705)
    img2_top = img1_top + img1_h + 0.50  # 1.65 + 2.20 + 0.50 = 4.35
    img2_w = LEFT_W - 0.20
    img2_h = img2_w / 2.705  # ≈ 1.57
    img2_left = LEFT_X + (LEFT_W - img2_w) / 2
    pic2 = slide.shapes.add_picture(
        str(RESOURCES / "国务院人工智能行动.png"),
        Inches(img2_left),
        Inches(img2_top),
        width=Inches(img2_w),
        height=Inches(img2_h),
    )
    pic2.line.color.rgb = COLOR_PRIMARY
    pic2.line.width = Pt(0.75)

    # Caption 2
    add_textbox(
        slide,
        left=LEFT_X,
        top=img2_top + img2_h + 0.05,
        width=LEFT_W,
        height=0.32,
        text='国务院 ·《关于深入实施"人工智能+"行动的意见》',
        size=11,
        bold=True,
        color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
    )

    # =========================================================
    # Right column: STAR narrative (S → T → A → R)
    # =========================================================
    RIGHT_X = 5.30
    RIGHT_W = 7.60

    # Helper to draw a STAR section
    def draw_star_section(top, label, label_color, title, body_lines, height=1.30):
        # Outer card background
        add_filled_rect(
            slide, RIGHT_X, top, RIGHT_W, height, fill=COLOR_GRAY_BG, corner=0.04
        )
        # Left vertical accent bar
        add_filled_rect(
            slide,
            RIGHT_X,
            top,
            0.18,
            height,
            fill=label_color,
            shape_type=MSO_SHAPE.RECTANGLE,
            corner=0.0,
        )
        # Big letter label (S/T/A/R)
        add_textbox(
            slide,
            left=RIGHT_X + 0.25,
            top=top + 0.05,
            width=0.65,
            height=height - 0.10,
            text=label,
            size=46,
            bold=True,
            color=label_color,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            font_name=EN_FONT,
        )
        # Section title
        add_textbox(
            slide,
            left=RIGHT_X + 1.00,
            top=top + 0.06,
            width=RIGHT_W - 1.10,
            height=0.34,
            text=title,
            size=14,
            bold=True,
            color=label_color,
        )
        # Body content
        body_box = slide.shapes.add_textbox(
            Inches(RIGHT_X + 1.00),
            Inches(top + 0.42),
            Inches(RIGHT_W - 1.10),
            Inches(height - 0.50),
        )
        add_paragraphs(body_box.text_frame, body_lines)

    # ---- S — Situation: 时代背景 ----
    s_body = [
        (
            "教育数字化已上升为国家战略 — 教育部《教育数字化战略行动》、"
            "国务院《人工智能+》行动意见明确要求推动 AI 与教育教学深度融合，"
            "智慧教育新形态加速到来。",
            {"size": 13, "color": COLOR_TEXT, "line_spacing": 1.45, "space_after": 6},
        ),
        (
            "高校实训教学是计算机类专业人才培养的关键环节，但当前的实训教学方式"
            "智能化程度普遍偏低、与产业需求衔接不足、教学资源更新缓慢，"
            "亟需借助大模型与 RAG 技术实现深度变革。",
            {"size": 13, "color": COLOR_TEXT, "line_spacing": 1.45},
        ),
    ]
    draw_star_section(
        top=1.20,
        label="S",
        label_color=COLOR_ACCENT,
        title="时代背景  ·  AI + 教育上升为国家战略",
        body_lines=s_body,
        height=2.55,
    )

    # ---- T — Task / Pain: 行业痛点 ----
    pain_color = RGBColor(0xC0, 0x39, 0x2B)  # warm red
    t_body = [
        (
            "❗  备课负担重     —     章节素材整理、题目设计动辄耗时数小时",
            {"size": 13, "color": COLOR_TEXT, "space_after": 8},
        ),
        (
            "❗  题库更新慢     —     考核内容陈旧、难以紧跟课程与产业演进",
            {"size": 13, "color": COLOR_TEXT, "space_after": 8},
        ),
        (
            "❗  个性化缺位     —     一对多教学难以兼顾学生差异化需求",
            {"size": 13, "color": COLOR_TEXT, "space_after": 8},
        ),
        (
            "❗  数据成孤岛     —     教与学过程数据分散、难以形成反馈闭环",
            {"size": 13, "color": COLOR_TEXT},
        ),
    ]
    draw_star_section(
        top=4.05,
        label="T",
        label_color=pain_color,
        title="行业痛点  ·  实训教学的四大顽疾",
        body_lines=t_body,
        height=2.95,
    )

    # =========================================================
    # Bottom highlight banner — grab attention
    # =========================================================
    banner_top = 7.10
    add_filled_rect(
        slide,
        0.55,
        banner_top,
        11.20,
        0.34,
        fill=COLOR_PRIMARY,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        corner=0.40,
    )
    add_textbox(
        slide,
        left=0.55,
        top=banner_top,
        width=11.20,
        height=0.34,
        text='💡  当政策春风遇上技术成熟期，正是智能教学辅助系统落地的最佳窗口',
        size=12,
        bold=True,
        color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_page_number(slide, 2)


def build_p3_objective(slide) -> None:
    """Page 3: 研究目标 — 针对四大痛点逐一给出解决对策（上下对比布局）."""
    prepare_content_slide(slide)
    add_title_bar(slide, "研究目标 — 针对痛点的系统性对策", "二")

    # =========================================================
    # Top banner: 总体研究目标
    # =========================================================
    add_filled_rect(slide, 0.55, 1.20, 12.30, 0.85, fill=COLOR_PRIMARY)
    add_textbox(
        slide,
        left=0.85,
        top=1.20,
        width=11.7,
        height=0.85,
        text="构建教师 / 学生 / 管理三端一体的 AI 教学辅助系统，"
             "针对四大痛点提供 RAG · 知识追踪 · 个性化推荐三位一体的系统性对策",
        size=14,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # =========================================================
    # Pain → Solution 4-column comparison
    # =========================================================
    pain_color = RGBColor(0xC0, 0x39, 0x2B)  # warm red
    solu_color = RGBColor(0x2E, 0x9C, 0x6E)  # green

    # Header labels
    add_textbox(
        slide,
        left=0.55,
        top=2.20,
        width=2.2,
        height=0.35,
        text="❗  痛点",
        size=14,
        bold=True,
        color=pain_color,
    )
    add_textbox(
        slide,
        left=0.55,
        top=4.55,
        width=2.2,
        height=0.35,
        text="✓  对策",
        size=14,
        bold=True,
        color=solu_color,
    )

    pairs = [
        {
            "icon_p": "📚",
            "pain_title": "备课负担重",
            "pain_desc": "章节素材整理、题目设计耗时数小时",
            "icon_s": "📝",
            "solu_title": "教师备课辅助",
            "solu_desc": "章节讲解提纲自动生成\n包含教学目标、重难点、课堂流程、实训任务",
        },
        {
            "icon_p": "📋",
            "pain_title": "题库更新慢",
            "pain_desc": "考核内容陈旧、与课程内容脱节",
            "icon_s": "✏️",
            "solu_title": "智能出题",
            "solu_desc": "按知识点定向生成 4 种题型\n单选 / 判断 / 填空 / 简答 + 标准答案",
        },
        {
            "icon_p": "🎯",
            "pain_title": "个性化缺位",
            "pain_desc": "一对多教学难以兼顾差异化需求",
            "icon_s": "🧭",
            "solu_title": "知识追踪 + 推荐",
            "solu_desc": "EMA 动态更新掌握度\n识别薄弱知识点并推送针对性练习",
        },
        {
            "icon_p": "🧩",
            "pain_title": "数据成孤岛",
            "pain_desc": "教与学过程数据分散、难以闭环",
            "icon_s": "🔄",
            "solu_title": "数据闭环",
            "solu_desc": "练习 → 评估 → 推荐\n作答数据自动驱动 EMA 更新",
        },
    ]

    # Layout config
    col_w = 2.85
    col_gap = 0.22
    cols_total_w = 4 * col_w + 3 * col_gap  # 12.06
    start_left = 0.55 + (12.30 - cols_total_w) / 2

    pain_top = 2.55
    pain_h = 1.65
    arrow_top = 4.25
    arrow_h = 0.30
    solu_top = 4.85
    solu_h = 2.05

    for i, p in enumerate(pairs):
        x = start_left + i * (col_w + col_gap)

        # ============ Pain card (top row) ============
        add_filled_rect(slide, x, pain_top, col_w, pain_h, fill=COLOR_LIGHT_BG, corner=0.06)
        # Top color band
        add_filled_rect(slide, x, pain_top, col_w, 0.50, fill=pain_color, corner=0.06)
        # Icon + title
        add_textbox(
            slide,
            left=x + 0.10,
            top=pain_top + 0.05,
            width=col_w - 0.20,
            height=0.40,
            text=f"{p['icon_p']}  {p['pain_title']}",
            size=14,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Pain description
        add_textbox(
            slide,
            left=x + 0.15,
            top=pain_top + 0.65,
            width=col_w - 0.30,
            height=pain_h - 0.75,
            text=p["pain_desc"],
            size=12,
            color=COLOR_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.4,
        )

        # ============ Down arrow ============
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.DOWN_ARROW,
            Inches(x + col_w / 2 - 0.18),
            Inches(arrow_top),
            Inches(0.36),
            Inches(arrow_h),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = COLOR_PRIMARY
        arrow.line.fill.background()

        # ============ Solution card (bottom row) ============
        add_filled_rect(slide, x, solu_top, col_w, solu_h, fill=COLOR_LIGHT_BG, corner=0.06)
        # Top color band
        add_filled_rect(slide, x, solu_top, col_w, 0.50, fill=solu_color, corner=0.06)
        add_textbox(
            slide,
            left=x + 0.10,
            top=solu_top + 0.05,
            width=col_w - 0.20,
            height=0.40,
            text=f"{p['icon_s']}  {p['solu_title']}",
            size=14,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Solution description
        add_textbox(
            slide,
            left=x + 0.15,
            top=solu_top + 0.65,
            width=col_w - 0.30,
            height=solu_h - 0.75,
            text=p["solu_desc"],
            size=11,
            color=COLOR_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.4,
        )

    # =========================================================
    # Bottom banner
    # =========================================================
    add_filled_rect(
        slide,
        0.55,
        7.10,
        11.20,
        0.34,
        fill=COLOR_PRIMARY,
        shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
        corner=0.40,
    )
    add_textbox(
        slide,
        left=0.55,
        top=7.10,
        width=11.20,
        height=0.34,
        text='🎯  痛点逐一被对策"接住" — 形成可量化、可落地、可持续演进的智能教学闭环',
        size=12,
        bold=True,
        color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_page_number(slide, 3)


def build_p4_use_case(slide) -> None:
    """Page 4: 用户角色与功能用例."""
    prepare_content_slide(slide)
    add_title_bar(slide, "用户角色与功能用例", "三")

    roles = [
        {
            "title": "教师角色",
            "subtitle": "Teacher",
            "color": COLOR_ACCENT,
            "items": [
                "UC-T1  上传课程知识库",
                "UC-T2  管理课程知识点",
                "UC-T3  生成章节讲解提纲",
                "UC-T4  按知识点智能出题",
                "UC-T5  查看学情分析",
            ],
        },
        {
            "title": "学生角色",
            "subtitle": "Student",
            "color": RGBColor(0x2E, 0x9C, 0x6E),
            "items": [
                "UC-S1  RAG 多轮问答",
                "UC-S2  完成在线练习",
                "UC-S3  查看知识掌握度",
                "UC-S4  接收个性化推荐",
                "UC-S5  查看作答历史",
            ],
        },
        {
            "title": "管理员",
            "subtitle": "Admin",
            "color": RGBColor(0xD9, 0x82, 0x2E),
            "items": [
                "UC-A1  用户账号管理",
                "UC-A2  课程资源管理",
                "UC-A3  数据看板查看",
                "UC-A4  权限分配",
                "UC-A5  系统配置",
            ],
        },
    ]
    card_top = 1.30
    card_h = 4.50
    card_w = 4.0
    gap = 0.18
    start_left = (SLIDE_W_IN - 3 * card_w - 2 * gap) / 2

    for i, role in enumerate(roles):
        left = start_left + i * (card_w + gap)
        # Card background
        add_filled_rect(slide, left, card_top, card_w, card_h, fill=COLOR_GRAY_BG, corner=0.05)
        # Header
        add_filled_rect(
            slide, left, card_top, card_w, 1.10, fill=role["color"], corner=0.05
        )
        add_textbox(
            slide,
            left=left,
            top=card_top + 0.20,
            width=card_w,
            height=0.45,
            text=role["title"],
            size=22,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide,
            left=left,
            top=card_top + 0.65,
            width=card_w,
            height=0.40,
            text=role["subtitle"],
            size=12,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=EN_FONT,
        )
        # Items
        box = slide.shapes.add_textbox(
            Inches(left + 0.30), Inches(card_top + 1.30), Inches(card_w - 0.60), Inches(card_h - 1.40)
        )
        add_paragraphs(
            box.text_frame,
            [(t, {"size": 13, "space_after": 10, "color": COLOR_TEXT}) for t in role["items"]],
        )

    # Bottom note
    add_textbox(
        slide,
        left=0.55,
        top=6.00,
        width=12.30,
        height=0.45,
        text="说明：系统按 RBAC（基于角色的访问控制）划分三类角色，所有角色共享统一的认证与会话管理。",
        size=12,
        color=COLOR_SUBTLE,
        align=PP_ALIGN.CENTER,
    )

    # Image placeholder for use-case diagram (optional)
    add_image_placeholder(
        slide,
        left=0.55,
        top=6.50,
        width=12.30,
        height=0.55,
        label="（可选）此处插入完整 UML 用例图（mermaid 渲染 PNG）",
    )

    add_page_number(slide, 4)


def build_p5_architecture(slide) -> None:
    """Page 5: 系统总体架构 + 技术栈（合并）.

    A 6-tier architecture diagram with technology stack tags annotated on
    each tier — combines the previous architecture and tech-stack pages
    into one comprehensive view.
    """
    prepare_content_slide(slide)
    add_title_bar(slide, "系统总体架构与技术栈", "四")

    # =========================================================
    # Top sub-title row
    # =========================================================
    add_textbox(
        slide,
        left=0.55,
        top=1.10,
        width=12.30,
        height=0.30,
        text="◆  6 层模块化架构  ·  前后端分离  ·  模型可替换  ·  流式优先  ·  数据闭环",
        size=13,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # =========================================================
    # 6-tier architecture stack
    # Each tier: [name | subtitle | content | tech tags]
    # =========================================================
    tiers = [
        {
            "name": "用户层",
            "name_en": "User",
            "items": "教师 · 学生 · 管理员",
            "tags": ["浏览器（Chrome / Edge / Safari）", "三角色 RBAC 权限"],
            "color": RGBColor(0x37, 0x4F, 0x7E),
        },
        {
            "name": "展示层",
            "name_en": "Presentation",
            "items": "教师端 · 学生端 · 管理端 — 11 个 Vue 页面",
            "tags": ["Vue 3", "Vite", "Element Plus", "Lucide Icons", "ECharts"],
            "color": RGBColor(0x4A, 0x90, 0xE2),
        },
        {
            "name": "业务层",
            "name_en": "Business",
            "items": "auth · courses · knowledge_base · rag_qa · exercises · lesson_plans · knowledge_tracking · conversations",
            "tags": ["FastAPI 0.110", "Pydantic", "uvicorn ASGI", "SSE 流式"],
            "color": RGBColor(0x2E, 0x9C, 0x6E),
        },
        {
            "name": "编排层",
            "name_en": "Orchestration",
            "items": "RAG Pipeline · 智能出题 · EMA 知识追踪 · 多轮对话记忆",
            "tags": ["LangChain", "Prompt 工程", "JSON Schema", "Few-shot"],
            "color": RGBColor(0xE5, 0xA5, 0x3B),
        },
        {
            "name": "AI 服务层",
            "name_en": "AI Services",
            "items": "大模型推理 · Embedding · Rerank · RAGAS 评测",
            "tags": ["Qwen-Plus / Turbo", "text-embedding-v3", "gte-rerank", "RAGAS"],
            "color": RGBColor(0xC0, 0x55, 0x55),
        },
        {
            "name": "数据层",
            "name_en": "Data",
            "items": "业务数据（用户/课程/对话/作答） · 向量索引 · 文档原文与切片",
            "tags": ["SQLite (8 表)", "ChromaDB (HNSW)", "BM25 + jieba", "本地文件"],
            "color": RGBColor(0x69, 0x40, 0x90),
        },
    ]

    base_top = 1.45
    tier_h = 0.86
    gap = 0.05
    layer_w = 12.30

    for i, t in enumerate(tiers):
        top = base_top + i * (tier_h + gap)

        # Tier background card
        add_filled_rect(slide, 0.55, top, layer_w, tier_h, fill=COLOR_LIGHT_BG, corner=0.04)
        # Left color label band
        add_filled_rect(slide, 0.55, top, 1.40, tier_h, fill=t["color"], corner=0.04)
        # Tier name (Chinese)
        add_textbox(
            slide,
            left=0.55,
            top=top + 0.08,
            width=1.40,
            height=0.42,
            text=t["name"],
            size=16,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        # Tier name (English)
        add_textbox(
            slide,
            left=0.55,
            top=top + 0.49,
            width=1.40,
            height=0.32,
            text=t["name_en"],
            size=9,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=EN_FONT,
        )
        # Tier content (left aligned)
        add_textbox(
            slide,
            left=2.05,
            top=top + 0.05,
            width=layer_w - 1.65,
            height=0.40,
            text=t["items"],
            size=11,
            color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Tech tags row (right side, like badges)
        tags_text = "    ".join(f"❑ {tag}" for tag in t["tags"])
        add_textbox(
            slide,
            left=2.05,
            top=top + 0.45,
            width=layer_w - 1.65,
            height=0.36,
            text=tags_text,
            size=10,
            bold=True,
            color=t["color"],
            anchor=MSO_ANCHOR.MIDDLE,
            font_name=EN_FONT,
        )

    # =========================================================
    # Down-arrow connectors between tiers (subtle)
    # =========================================================
    for i in range(len(tiers) - 1):
        y = base_top + (i + 1) * tier_h + i * gap + 0.01
        # tiny down-arrow at left & right edges
        for x_pos in [3.5, 9.5]:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(x_pos),
                Inches(y),
                Inches(0.18),
                Inches(0.10),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_SUBTLE
            arrow.line.fill.background()

    # =========================================================
    # Bottom: data flow note
    # =========================================================
    add_textbox(
        slide,
        left=0.55,
        top=7.10,
        width=12.30,
        height=0.32,
        text="🔄  数据闭环：学生作答 → exercise_attempts → EMA 更新 knowledge_mastery → 识别薄弱点 → 个性化推荐",
        size=11,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )

    add_page_number(slide, 5)


def build_p7_database(slide) -> None:
    """Page 7 (now 6): 核心数据库设计 — 8 表 + ER 占位 + 数据闭环."""
    prepare_content_slide(slide)
    add_title_bar(slide, "核心数据库设计", "五")

    # =========================================================
    # Section 1: 8 tables grid (compact 2×4)
    # =========================================================
    tables = [
        ("users", "用户账号\n三角色", COLOR_ACCENT),
        ("sessions", "登录会话\ntoken 管理", COLOR_ACCENT),
        ("courses", "课程信息\n基本元数据", COLOR_ACCENT),
        ("knowledge_points", "课程知识点\n自动+手动", COLOR_ACCENT),
        ("conversations", "对话会话\n多轮问答", RGBColor(0x2E, 0x9C, 0x6E)),
        ("messages", "对话消息\n含引用 JSON", RGBColor(0x2E, 0x9C, 0x6E)),
        ("knowledge_mastery", "EMA 掌握度 ⭐\n知识追踪", RGBColor(0xD9, 0x82, 0x2E)),
        ("exercise_attempts", "作答记录 ⭐\n闭环数据", RGBColor(0xD9, 0x82, 0x2E)),
    ]

    # Section title
    add_textbox(
        slide,
        left=0.55,
        top=1.10,
        width=12.30,
        height=0.32,
        text='◆  8 张核心业务表  —  按"基础数据 / 对话 / 知识追踪"三组着色',
        size=12,
        bold=True,
        color=COLOR_PRIMARY,
    )

    cell_w = 2.90
    cell_h = 1.05
    gap_x = 0.20
    gap_y = 0.15
    start_left = (SLIDE_W_IN - 4 * cell_w - 3 * gap_x) / 2
    base_top = 1.45
    for idx, (name, desc, color) in enumerate(tables):
        col = idx % 4
        row = idx // 4
        left = start_left + col * (cell_w + gap_x)
        top = base_top + row * (cell_h + gap_y)
        add_filled_rect(slide, left, top, cell_w, cell_h, fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, left, top, cell_w, 0.36, fill=color, corner=0.06)
        add_textbox(
            slide,
            left=left,
            top=top + 0.02,
            width=cell_w,
            height=0.34,
            text=name,
            size=12,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            font_name=EN_FONT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left=left + 0.08,
            top=top + 0.40,
            width=cell_w - 0.16,
            height=cell_h - 0.45,
            text=desc,
            size=11,
            color=COLOR_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # =========================================================
    # Section 2: 左 ER 关系图占位符 + 右 数据闭环
    # =========================================================
    section_top = 3.85
    section_h = 3.30

    # ---- Left: ER diagram (auto-generated from backend/app/db.py) ----
    er_left = 0.55
    er_w = 5.40
    add_filled_rect(slide, er_left, section_top, er_w, 0.40, fill=COLOR_PRIMARY)
    add_textbox(
        slide,
        left=er_left,
        top=section_top,
        width=er_w,
        height=0.40,
        text="◆  ER 关系示意图（自动生成自 db.py）",
        size=13,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    # Light background card
    add_filled_rect(
        slide,
        er_left,
        section_top + 0.40,
        er_w,
        section_h - 0.40,
        fill=RGBColor(0xFA, 0xFC, 0xFE),
        corner=0.04,
    )
    # Embed ER diagram image (centered, scaled to fit container)
    if ER_DIAGRAM_PNG.exists():
        # Container dims
        cx_left = er_left + 0.05
        cx_top = section_top + 0.45
        cx_w = er_w - 0.10
        cx_h = section_h - 0.50

        # Inspect actual image dims to compute aspect-correct fit
        try:
            from PIL import Image as _PILImage  # type: ignore

            with _PILImage.open(ER_DIAGRAM_PNG) as im:
                iw_px, ih_px = im.size
            img_ratio = iw_px / ih_px
        except Exception:
            img_ratio = 0.83  # fallback (matches the current dot output)

        cont_ratio = cx_w / cx_h
        if img_ratio >= cont_ratio:
            # Wider than container: fit by width
            draw_w = cx_w
            draw_h = draw_w / img_ratio
        else:
            # Taller than container: fit by height
            draw_h = cx_h
            draw_w = draw_h * img_ratio
        draw_left = cx_left + (cx_w - draw_w) / 2
        draw_top = cx_top + (cx_h - draw_h) / 2
        slide.shapes.add_picture(
            str(ER_DIAGRAM_PNG),
            Inches(draw_left),
            Inches(draw_top),
            width=Inches(draw_w),
            height=Inches(draw_h),
        )
    else:
        add_textbox(
            slide,
            left=er_left + 0.20,
            top=section_top + 0.60,
            width=er_w - 0.40,
            height=section_h - 0.80,
            text=(
                "📷  ER 图未生成\n\n"
                "请先运行：python scripts/export_er_diagram.py\n"
                "脚本会从 backend/app/db.py 自动生成图片到\n"
                "data/er-diagram/er.png"
            ),
            size=11,
            color=RGBColor(0x66, 0x66, 0x66),
            line_spacing=1.4,
        )

    # ---- Right: data closed-loop ----
    loop_left = 6.20
    loop_w = 6.65
    add_filled_rect(slide, loop_left, section_top, loop_w, 0.40, fill=COLOR_ACCENT)
    add_textbox(
        slide,
        left=loop_left,
        top=section_top,
        width=loop_w,
        height=0.40,
        text="🔄  数据闭环：练习 → 评估 → 推荐",
        size=13,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    add_filled_rect(
        slide,
        loop_left,
        section_top + 0.40,
        loop_w,
        section_h - 0.40,
        fill=COLOR_GRAY_BG,
        corner=0.04,
    )

    flow_steps = [
        ("学生作答", "exercise_\nattempts 写入"),
        ("EMA 算法", "更新 knowledge\n_mastery"),
        ("识别薄弱", "掌握度 < 0.6\n的知识点"),
        ("Prompt 注入", "组装个性化\n推荐 Prompt"),
        ("LLM 生成", "Qwen-Plus 生成\n针对性练习"),
        ("推送学生", "前端展示\n推荐题目"),
    ]
    # 2 rows × 3 cols
    sb_w = 1.95
    sb_h = 1.20
    sb_gx = 0.10
    sb_gy = 0.15
    sb_start = loop_left + (loop_w - 3 * sb_w - 2 * sb_gx) / 2
    sb_top = section_top + 0.55

    for i, (head, body) in enumerate(flow_steps):
        col = i % 3
        row = i // 3
        x = sb_start + col * (sb_w + sb_gx)
        y = sb_top + row * (sb_h + sb_gy)
        add_filled_rect(slide, x, y, sb_w, sb_h, fill=RGBColor(0xFF, 0xFF, 0xFF), corner=0.06)
        add_filled_rect(slide, x, y, sb_w, 0.34, fill=COLOR_PRIMARY, corner=0.06)
        add_textbox(
            slide,
            left=x,
            top=y + 0.02,
            width=sb_w,
            height=0.32,
            text=f"{i+1}. {head}",
            size=11,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left=x + 0.05,
            top=y + 0.36,
            width=sb_w - 0.10,
            height=sb_h - 0.40,
            text=body,
            size=10,
            color=COLOR_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
            font_name=EN_FONT,
        )

    add_page_number(slide, 6)


def build_p8_teacher_overview(slide) -> None:
    """Page 8: 教师侧功能总览."""
    prepare_content_slide(slide)
    add_title_bar(slide, "教师侧功能总览", "六")

    # Top: task book mapping table
    add_textbox(
        slide,
        left=0.55,
        top=1.30,
        width=12.30,
        height=0.45,
        text="◆ 任务书要求 → 系统实现映射",
        size=16,
        bold=True,
        color=COLOR_PRIMARY,
    )
    rows = [
        ("基于课程大纲和本地知识库自动生成章节知识讲解提纲", "模块 2：教师备课辅助", "✅"),
        ("基于课程大纲自动生成基础实训练习", "模块 2 输出 \"基础实训任务\" 字段", "✅"),
        ("按知识点自动生成练习题与测试题", "模块 3：智能出题", "✅"),
        ("提供参考答案", "模块 3：每题附标准答案 + 评分标准", "✅"),
    ]
    table_top = 1.80
    # Header
    add_filled_rect(slide, 0.55, table_top, 12.30, 0.50, fill=COLOR_PRIMARY, corner=0.02)
    headers = ["任务书要求", "系统实现模块", "状态"]
    widths = [6.50, 4.60, 1.20]
    xs = [0.55, 7.05, 11.65]
    for h, x, w in zip(headers, xs, widths):
        add_textbox(
            slide,
            left=x,
            top=table_top,
            width=w,
            height=0.50,
            text=h,
            size=14,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    # Data rows
    for ri, row in enumerate(rows):
        row_top = table_top + 0.50 + ri * 0.45
        if ri % 2 == 0:
            add_filled_rect(slide, 0.55, row_top, 12.30, 0.45, fill=COLOR_LIGHT_BG, corner=0.0)
        for c, x, w in zip(row, xs, widths):
            add_textbox(
                slide,
                left=x + 0.10,
                top=row_top,
                width=w - 0.20,
                height=0.45,
                text=c,
                size=12,
                color=COLOR_TEXT if c != "✅" else RGBColor(0x2E, 0x9C, 0x6E),
                bold=c == "✅",
                align=PP_ALIGN.CENTER if c == "✅" else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # Bottom: 3 module cards
    add_textbox(
        slide,
        left=0.55,
        top=4.40,
        width=12.30,
        height=0.40,
        text="◆ 三大核心模块（详见后续 3 页）",
        size=16,
        bold=True,
        color=COLOR_PRIMARY,
    )
    modules = [
        {
            "no": "①",
            "title": "知识库管理 + RAG 问答",
            "desc": "上传 PDF / Word / Markdown / 网页 URL\n自动解析切分 → 向量化入库\n混合检索 + 重排序 + 流式问答",
        },
        {
            "no": "②",
            "title": "教师备课辅助",
            "desc": "输入：课程 + 章节 + 课时 + 知识点\n输出：6 大模块结构化讲解提纲\n包含教学目标 / 重难点 / 课堂流程",
        },
        {
            "no": "③",
            "title": "智能出题",
            "desc": "题型：单选 / 判断 / 填空 / 简答\n按章节 + 知识点 + 难度三维定向\n标准答案 + 知识点标签 + 评分标准",
        },
    ]
    card_top = 4.90
    card_h = 2.20
    card_w = 4.0
    gap = 0.15
    start_left = (SLIDE_W_IN - 3 * card_w - 2 * gap) / 2
    for i, m in enumerate(modules):
        left = start_left + i * (card_w + gap)
        add_filled_rect(slide, left, card_top, card_w, card_h, fill=COLOR_LIGHT_BG, corner=0.05)
        # Number
        add_textbox(
            slide,
            left=left + 0.20,
            top=card_top + 0.10,
            width=0.60,
            height=0.60,
            text=m["no"],
            size=32,
            bold=True,
            color=COLOR_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Title
        add_textbox(
            slide,
            left=left + 0.85,
            top=card_top + 0.20,
            width=card_w - 0.95,
            height=0.45,
            text=m["title"],
            size=15,
            bold=True,
            color=COLOR_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Description
        add_textbox(
            slide,
            left=left + 0.20,
            top=card_top + 0.85,
            width=card_w - 0.40,
            height=card_h - 1.00,
            text=m["desc"],
            size=11,
            color=COLOR_TEXT,
            line_spacing=1.4,
        )

    add_page_number(slide, 7)


def build_p9_rag(slide) -> None:
    """Page 9: 模块 1 — 知识库管理与 RAG 问答."""
    prepare_content_slide(slide)
    add_title_bar(slide, "模块 1 — 知识库管理与 RAG 问答", "七")

    # Top: pipeline (offline + online)
    add_textbox(
        slide,
        left=0.55,
        top=1.25,
        width=12.30,
        height=0.40,
        text="◆ RAG Pipeline：离线知识库构建 + 在线问答检索",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # Offline pipeline
    add_textbox(
        slide,
        left=0.55,
        top=1.70,
        width=12.30,
        height=0.30,
        text="离线阶段：",
        size=11,
        bold=True,
        color=COLOR_ACCENT,
    )
    offline_steps = [
        "文档上传\nPDF/Word/MD/URL",
        "解析\nPyPDF/docx",
        "切分\n规则/LLM 辅助",
        "jieba 分词",
        "Embedding\ntext-embed-v3",
        "ChromaDB\nHNSW 索引",
    ]
    _draw_pipeline(slide, offline_steps, top=2.05, color=RGBColor(0xE5, 0xA5, 0x3B))

    # Online pipeline
    add_textbox(
        slide,
        left=0.55,
        top=3.10,
        width=12.30,
        height=0.30,
        text="在线阶段：",
        size=11,
        bold=True,
        color=COLOR_ACCENT,
    )
    online_steps = [
        "用户提问",
        "BM25 + 向量\n并行召回",
        "融合去重",
        "gte-rerank\n重排序",
        "Prompt 拼装\n+ 历史 + 联网",
        "Qwen-Plus\n流式生成",
        "SSE 推送 +\n引用回显",
    ]
    _draw_pipeline(slide, online_steps, top=3.45, color=COLOR_ACCENT)

    # Compact features row (one-liner) — keeps page tighter
    add_textbox(
        slide,
        left=0.55,
        top=4.50,
        width=12.30,
        height=0.28,
        text="◆ 核心特性： ✓ BM25+向量混合检索   ✓ Cross-Encoder 重排   ✓ 滑窗 10 条多轮记忆   ✓ DashScope 联网搜索   ✓ SSE 流式 (首字 <1s)   ✓ RAGAS 评测",
        size=11,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # =========================================================
    # 3 large screenshot placeholders (with capture guide)
    # =========================================================
    add_textbox(
        slide,
        left=0.55,
        top=4.85,
        width=12.30,
        height=0.28,
        text="◆ 系统演示截图（待截取）",
        size=13,
        bold=True,
        color=COLOR_ACCENT,
    )

    sc_top = 5.18
    sc_h = 1.45
    sc_gap = 0.18
    sc_w = (12.30 - 2 * sc_gap) / 3  # ≈ 4.0

    placeholders = [
        (
            "截图 1：知识库上传与文档管理",
            "页面 /knowledge-base  ·  KnowledgeBaseUpload.vue\n"
            "上传 1 份 PDF（如《大模型应用开发》）→\n"
            '截"已上传文档列表 + 切片数 + 处理状态"区域',
        ),
        (
            "截图 2：知识库检索与召回结果",
            "页面 /knowledge-base/search  ·  KnowledgeBaseSearch.vue\n"
            '输入查询"什么是 RAG" → 截显示混合检索结果 +\n'
            "rerank 分数 + 来源文档片段的列表",
        ),
        (
            "截图 3：RAG 多轮问答（含引用）",
            "页面 /qa  ·  RagQa.vue\n"
            '提问"对比 BM25 与向量检索" → 截带引用'
            "[1][2] 标注 + 联网来源 + 流式生成中的对话气泡",
        ),
    ]
    for i, (title, guide) in enumerate(placeholders):
        x = 0.55 + i * (sc_w + sc_gap)
        # Image placeholder
        add_image_placeholder(
            slide,
            left=x,
            top=sc_top,
            width=sc_w,
            height=sc_h,
            label=title,
        )
        # Capture-guide caption underneath
        add_textbox(
            slide,
            left=x,
            top=sc_top + sc_h + 0.04,
            width=sc_w,
            height=0.75,
            text=guide,
            size=9,
            color=RGBColor(0x55, 0x55, 0x55),
            line_spacing=1.20,
        )

    add_page_number(slide, 8)


def _draw_pipeline(slide, steps: list[str], top: float, color: RGBColor) -> None:
    """Draw a horizontal pipeline of boxes connected by arrows."""
    n = len(steps)
    arrow_w = 0.18
    total_arrow = (n - 1) * arrow_w
    box_w = (12.30 - total_arrow) / n
    box_h = 0.95
    x = 0.55
    for i, label in enumerate(steps):
        add_filled_rect(slide, x, top, box_w, box_h, fill=COLOR_LIGHT_BG, corner=0.06)
        # Top color strip
        add_filled_rect(slide, x, top, box_w, 0.20, fill=color, corner=0.06)
        add_textbox(
            slide,
            left=x + 0.05,
            top=top + 0.22,
            width=box_w - 0.10,
            height=box_h - 0.25,
            text=label,
            size=10,
            color=COLOR_TEXT,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        if i < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_w),
                Inches(top + 0.35),
                Inches(arrow_w),
                Inches(0.25),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = color
            arrow.line.fill.background()
        x += box_w + arrow_w


def build_p10_lesson_plan(slide) -> None:
    """Page 10: 模块 2 — 教师备课辅助."""
    prepare_content_slide(slide)
    add_title_bar(slide, "模块 2 — 教师备课辅助（章节讲解提纲）", "八")

    # =========================================================
    # Left column (4.20 wide): Input + flow
    # =========================================================
    # ---- Input parameters card ----
    add_filled_rect(slide, 0.55, 1.20, 4.20, 2.60, fill=COLOR_LIGHT_BG, corner=0.06)
    add_filled_rect(slide, 0.55, 1.20, 4.20, 0.45, fill=COLOR_PRIMARY, corner=0.06)
    add_textbox(
        slide,
        left=0.55,
        top=1.20,
        width=4.20,
        height=0.45,
        text="📥  输入参数",
        size=14,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    inputs = [
        ("课程：", "大模型应用开发"),
        ("章节：", "第 3 章 RAG 检索增强生成"),
        ("课时：", "90 分钟"),
        ("知识点：", "向量检索 / BM25 / Rerank"),
    ]
    box = slide.shapes.add_textbox(Inches(0.75), Inches(1.80), Inches(3.85), Inches(1.95))
    items = []
    for k, v in inputs:
        items.append((k, {"size": 11, "bold": True, "color": COLOR_PRIMARY, "space_after": 2}))
        items.append((f"  {v}", {"size": 11, "color": COLOR_TEXT, "space_after": 6}))
    add_paragraphs(box.text_frame, items)

    # ---- Generation flow ----
    add_filled_rect(slide, 0.55, 3.95, 4.20, 1.50, fill=COLOR_GRAY_BG, corner=0.06)
    add_textbox(
        slide,
        left=0.55,
        top=4.00,
        width=4.20,
        height=0.30,
        text="◆ 生成流程",
        size=12,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    flow_items = [
        "① 检索知识库 Top-N 片段",
        "② 组装 Prompt（JSON Schema + Few-shot）",
        "③ Qwen-Plus 生成",
        "④ 6 大模块结构化输出",
        "⑤ 教师审核 / 微调 / 导出",
    ]
    box = slide.shapes.add_textbox(Inches(0.75), Inches(4.30), Inches(3.85), Inches(1.10))
    add_paragraphs(
        box.text_frame,
        [(s, {"size": 10, "space_after": 2, "color": COLOR_TEXT}) for s in flow_items],
    )

    # ---- Screenshot placeholder 1 (input form) ----
    add_image_placeholder(
        slide,
        left=0.55,
        top=5.60,
        width=4.20,
        height=1.45,
        label="截图 1：备课参数输入表单\nLessonOutline.vue（顶部）",
    )

    # =========================================================
    # Right column (8.10 wide): 6 output modules + screenshot
    # =========================================================
    outputs = [
        ("①", "教学目标", "掌握 RAG 基本原理；理解混合检索优势；能动手实现简单 RAG Demo"),
        ("②", "重点难点", "重点：检索召回与重排序  |  难点：Prompt 工程与幻觉抑制"),
        ("③", "课堂流程", "导入 10 → 原理讲解 30 → 演示 20 → 实训 25 → 总结 5（min）"),
        ("④", "基础实训任务", "任务 1：搭建 ChromaDB；任务 2：实现混合检索；任务 3：完成 Demo"),
        ("⑤", "考核建议", "课堂提问 3 题 + 实训代码提交 + 简答题 1 题"),
        ("⑥", "知识来源", "引用知识库片段 [1][2][3]，含文档名 + 页码 + 原文摘录"),
    ]
    base_top = 1.20
    item_h = 0.62
    item_w = 8.10
    item_left = 4.85
    for i, (no, title, desc) in enumerate(outputs):
        top = base_top + i * (item_h + 0.05)
        add_filled_rect(slide, item_left, top, item_w, item_h, fill=COLOR_GRAY_BG, corner=0.04)
        add_filled_rect(slide, item_left, top, 0.45, item_h, fill=COLOR_PRIMARY, corner=0.04)
        add_textbox(
            slide,
            left=item_left,
            top=top,
            width=0.45,
            height=item_h,
            text=no,
            size=16,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left=item_left + 0.55,
            top=top + 0.02,
            width=1.50,
            height=item_h - 0.04,
            text=title,
            size=11,
            bold=True,
            color=COLOR_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left=item_left + 2.10,
            top=top + 0.02,
            width=item_w - 2.20,
            height=item_h - 0.04,
            text=desc,
            size=10,
            color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ---- Screenshot placeholder 2 (output result) ----
    add_image_placeholder(
        slide,
        left=4.85,
        top=5.30,
        width=8.10,
        height=1.75,
        label="截图 2：6 大模块结构化讲解提纲（生成结果）— LessonOutline.vue 输出区",
    )

    # =========================================================
    # Bottom value banner
    # =========================================================
    add_textbox(
        slide,
        left=0.55,
        top=7.15,
        width=12.30,
        height=0.30,
        text='💡  备课时间：从"几小时整理素材" → "10 秒生成 + 微调"，输出可直接用于实训课堂',
        size=12,
        bold=True,
        color=RGBColor(0xD9, 0x82, 0x2E),
        align=PP_ALIGN.CENTER,
    )

    add_page_number(slide, 9)


def build_p11_exercise(slide) -> None:
    """Page 11: 模块 3 — 智能出题."""
    prepare_content_slide(slide)
    add_title_bar(slide, "模块 3 — 智能出题", "九")

    # Top: 4 question types
    add_textbox(
        slide,
        left=0.55,
        top=1.25,
        width=12.30,
        height=0.40,
        text="◆ 支持 4 种题型，全面覆盖任务书要求",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )
    types = [
        {
            "name": "单选题",
            "en": "single_choice",
            "fields": "stem + options[4] + answer",
            "eval": "选项精确匹配",
        },
        {
            "name": "判断题",
            "en": "true_false",
            "fields": "stem + answer(bool)",
            "eval": "布尔值匹配",
        },
        {
            "name": "填空题 ⭐",
            "en": "fill_in_blank",
            "fields": "stem + blanks[多空+替代答案]",
            "eval": "逐空 + alternatives + 忽略大小写",
        },
        {
            "name": "简答题",
            "en": "short_answer",
            "fields": "stem + reference_answer + rubric",
            "eval": "关键词 + 语义相似度（规划）",
        },
    ]
    card_w = 2.95
    gap = 0.14
    base_top = 1.75
    card_h = 1.65
    start_left = (SLIDE_W_IN - 4 * card_w - 3 * gap) / 2
    for i, t in enumerate(types):
        left = start_left + i * (card_w + gap)
        add_filled_rect(slide, left, base_top, card_w, card_h, fill=COLOR_LIGHT_BG, corner=0.06)
        add_filled_rect(slide, left, base_top, card_w, 0.45, fill=COLOR_PRIMARY, corner=0.06)
        add_textbox(
            slide,
            left=left,
            top=base_top + 0.02,
            width=card_w,
            height=0.43,
            text=t["name"],
            size=14,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide,
            left=left + 0.10,
            top=base_top + 0.50,
            width=card_w - 0.20,
            height=0.30,
            text=t["en"],
            size=10,
            color=COLOR_ACCENT,
            align=PP_ALIGN.CENTER,
            font_name=EN_FONT,
        )
        add_textbox(
            slide,
            left=left + 0.10,
            top=base_top + 0.80,
            width=card_w - 0.20,
            height=0.40,
            text=f"字段：{t['fields']}",
            size=10,
            color=COLOR_TEXT,
        )
        add_textbox(
            slide,
            left=left + 0.10,
            top=base_top + 1.20,
            width=card_w - 0.20,
            height=0.40,
            text=f"评测：{t['eval']}",
            size=10,
            color=COLOR_TEXT,
        )

    # Generation flow
    add_textbox(
        slide,
        left=0.55,
        top=3.55,
        width=12.30,
        height=0.40,
        text="◆ 生成流程",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )
    flow_steps = [
        "教师选参数\n课程·题型·数量·难度·知识点",
        "按知识点检索\n知识库 Top-N 片段",
        "组装 Prompt\nJSON Schema + Few-shot",
        "Qwen-Plus 调用",
        "后端校验 + 知识点回填",
        "持久化 +\n前端列表展示",
    ]
    _draw_pipeline(slide, flow_steps, top=4.00, color=COLOR_PRIMARY)

    # Compact prompt-engineering badges (one row)
    add_textbox(
        slide,
        left=0.55,
        top=5.00,
        width=12.30,
        height=0.28,
        text="◆ Prompt 工程：🔧 JSON Schema 结构化输出   🎯 Few-shot 示例   📊 难度指令   🎓 知识点定向   🛡️ 检索片段防幻觉",
        size=11,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # =========================================================
    # 3 large screenshot placeholders + capture guide
    # =========================================================
    add_textbox(
        slide,
        left=0.55,
        top=5.32,
        width=12.30,
        height=0.28,
        text="◆ 系统演示截图（待截取）",
        size=13,
        bold=True,
        color=COLOR_ACCENT,
    )

    sc_top = 5.65
    sc_h = 1.20
    sc_gap = 0.18
    sc_w = (12.30 - 2 * sc_gap) / 3

    placeholders = [
        (
            "截图 1：出题表单",
            "页面 /exercises  ·  ExerciseGeneration.vue\n"
            "选课程→题型(混合)→数量(10)→难度(中)→\n"
            '勾选 3 个知识点 → 截"参数表单"区',
        ),
        (
            "截图 2：生成结果列表",
            '点击"开始生成"按钮 →\n'
            "截 4 种题型混合的题目列表（题干 +\n"
            "题型徽章 + 难度标签）",
        ),
        (
            "截图 3：题目详情 + 答案",
            "点击列表中任一题 → 截弹窗 / 详情页：\n"
            "题干 + 选项 / 填空 + 标准答案 +\n"
            "知识点标签 + 评分标准",
        ),
    ]
    for i, (title, guide) in enumerate(placeholders):
        x = 0.55 + i * (sc_w + sc_gap)
        add_image_placeholder(
            slide,
            left=x,
            top=sc_top,
            width=sc_w,
            height=sc_h,
            label=title,
        )
        add_textbox(
            slide,
            left=x,
            top=sc_top + sc_h + 0.03,
            width=sc_w,
            height=0.60,
            text=guide,
            size=9,
            color=RGBColor(0x55, 0x55, 0x55),
            line_spacing=1.18,
        )

    add_page_number(slide, 10)


def build_p12_progress(slide) -> None:
    """Page 12: 进度对照与阶段性总结."""
    prepare_content_slide(slide)
    add_title_bar(slide, "进度对照与阶段性总结", "十")

    # Left 65%: progress comparison table
    add_textbox(
        slide,
        left=0.55,
        top=1.25,
        width=8.20,
        height=0.40,
        text="◆ 进度对照（对照任务书 16 周安排）",
        size=14,
        bold=True,
        color=COLOR_PRIMARY,
    )
    headers = ["周次", "任务书安排", "完成情况", "状态"]
    widths = [1.10, 3.50, 3.00, 0.60]
    xs = [0.55, 1.65, 5.15, 8.15]
    rows = [
        ("1-2 周", "文献调研 / 开题 / 外文翻译", "开题 + 综述 + 2 篇翻译", "✅"),
        ("3-5 周", "需求 / 概要 / 详细 / 数据库设计", "4 层架构 + 8 张表 + API", "✅"),
        ("6-8 周", "教师侧功能模块编码", "知识库 + 备课 + 出题完成", "✅"),
        ("第 9 周", "中期检查", "当前阶段", "🟡"),
        ("10-12 周", "学生侧 / 管理侧其余模块", "已提前完成知识追踪闭环", "🔵"),
        ("13-14 周", "测试 / 论文 / 查重", "待开始", "🔲"),
        ("15-16 周", "PPT / 预答辩 / 答辩", "待开始", "🔲"),
    ]
    table_top = 1.70
    # Header
    add_filled_rect(slide, 0.55, table_top, 8.20, 0.40, fill=COLOR_PRIMARY, corner=0.02)
    for h, x, w in zip(headers, xs, widths):
        add_textbox(
            slide,
            left=x,
            top=table_top,
            width=w,
            height=0.40,
            text=h,
            size=11,
            bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )
    # Data rows
    for ri, row in enumerate(rows):
        row_top = table_top + 0.40 + ri * 0.42
        if ri % 2 == 0:
            add_filled_rect(slide, 0.55, row_top, 8.20, 0.42, fill=COLOR_LIGHT_BG, corner=0.0)
        for c, x, w in zip(row, xs, widths):
            color = COLOR_TEXT
            bold = False
            if c == "🟡" or c == "🔵" or c == "✅":
                bold = True
            add_textbox(
                slide,
                left=x + 0.05,
                top=row_top,
                width=w - 0.10,
                height=0.42,
                text=c,
                size=10,
                bold=bold,
                color=color,
                align=PP_ALIGN.CENTER if w < 2 else PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE,
            )

    # Quantitative metrics row (4 cards)
    add_textbox(
        slide,
        left=0.55,
        top=4.85,
        width=8.20,
        height=0.40,
        text="◆ 完成度量化指标",
        size=14,
        bold=True,
        color=COLOR_PRIMARY,
    )
    metrics = [
        ("9", "Router 模块"),
        ("11", "Service 模块"),
        ("11", "Vue 页面"),
        ("8", "数据库表"),
        ("4", "题型支持"),
        ("4", "RAGAS 指标"),
    ]
    m_w = 1.25
    m_gap = 0.10
    m_top = 5.30
    m_h = 1.10
    m_start = 0.55
    for i, (num, label) in enumerate(metrics):
        x = m_start + i * (m_w + m_gap)
        add_filled_rect(slide, x, m_top, m_w, m_h, fill=COLOR_PRIMARY, corner=0.06)
        add_textbox(
            slide,
            left=x,
            top=m_top + 0.05,
            width=m_w,
            height=0.55,
            text=num,
            size=28,
            bold=True,
            color=RGBColor(0xFF, 0xCC, 0x00),
            align=PP_ALIGN.CENTER,
            font_name=EN_FONT,
        )
        add_textbox(
            slide,
            left=x,
            top=m_top + 0.65,
            width=m_w,
            height=0.40,
            text=label,
            size=10,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )

    # Right 35%: achievements
    add_filled_rect(slide, 9.05, 1.25, 3.85, 5.85, fill=COLOR_LIGHT_BG, corner=0.05)
    add_filled_rect(slide, 9.05, 1.25, 3.85, 0.50, fill=COLOR_PRIMARY, corner=0.05)
    add_textbox(
        slide,
        left=9.05,
        top=1.25,
        width=3.85,
        height=0.50,
        text="✦ 阶段性核心成果",
        size=14,
        bold=True,
        color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )
    items = [
        ("🟢 已完成", {"size": 12, "bold": True, "color": RGBColor(0x2E, 0x9C, 0x6E), "space_after": 2}),
        ("• 教师侧 3 大模块全部完成", {"size": 10, "space_after": 2}),
        ("• 任务书教师侧目标 100% 达成", {"size": 10, "space_after": 8}),
        ("🌟 超额完成（不在原计划内）", {"size": 12, "bold": True, "color": RGBColor(0xD9, 0x82, 0x2E), "space_after": 2}),
        ("• 混合检索 + Rerank", {"size": 10, "space_after": 2}),
        ("• 多轮对话记忆", {"size": 10, "space_after": 2}),
        ("• 联网搜索 + 角标引用", {"size": 10, "space_after": 2}),
        ("• 填空题题型", {"size": 10, "space_after": 2}),
        ("• 大模型辅助语义切分", {"size": 10, "space_after": 2}),
        ("• RAGAS 评测接口", {"size": 10, "space_after": 2}),
        ("• 知识追踪闭环（提前完成）", {"size": 10, "space_after": 8}),
        ("🟢 结论", {"size": 12, "bold": True, "color": COLOR_PRIMARY, "space_after": 2}),
        ("进度严格符合任务书要求且部分超前，可保质保量完成毕业设计。", {"size": 10, "color": COLOR_TEXT}),
    ]
    box = slide.shapes.add_textbox(Inches(9.20), Inches(1.85), Inches(3.55), Inches(5.20))
    add_paragraphs(box.text_frame, items)

    add_page_number(slide, 11)


def build_p13_remaining(slide) -> None:
    """Page 13: 剩余工作计划."""
    prepare_content_slide(slide)
    add_title_bar(slide, "剩余工作计划", "十一")

    # Top: Gantt-like timeline
    add_textbox(
        slide,
        left=0.55,
        top=1.25,
        width=12.30,
        height=0.40,
        text="◆ 后续 8 周工作计划（第 9 周 → 第 16 周）",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )

    # Week labels
    weeks = ["W9", "W10", "W11", "W12", "W13", "W14", "W15", "W16"]
    track_left = 3.10
    track_right = 12.85
    track_w = track_right - track_left
    week_w = track_w / 8
    week_top = 1.75
    for i, w in enumerate(weeks):
        x = track_left + i * week_w
        add_filled_rect(slide, x, week_top, week_w, 0.40, fill=COLOR_GRAY_BG, corner=0.0)
        add_textbox(
            slide,
            left=x,
            top=week_top,
            width=week_w,
            height=0.40,
            text=w,
            size=11,
            bold=True,
            color=COLOR_PRIMARY,
            align=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # Tasks (label + start_week (0-indexed) + duration in weeks)
    tasks = [
        ("① 主观题高级评分（语义相似度 + LLM 加权融合）", 1, 1, RGBColor(0x4A, 0x90, 0xE2)),
        ("② 学情分析（错题聚合 + 班级统计）", 2, 1, RGBColor(0x2E, 0x9C, 0x6E)),
        ("③ 管理端数据看板（ECharts 图表）", 2, 1, RGBColor(0xE5, 0xA5, 0x3B)),
        ("④ 系统稳定性完善（异常处理 + 性能优化）", 3, 1, RGBColor(0xC0, 0x55, 0x55)),
        ("⑤ 系统测试（功能 + 性能 + 用户测试）", 4, 1, RGBColor(0x9B, 0x5E, 0xC2)),
        ("⑥ 论文撰写 + 查重修改", 4, 2, RGBColor(0x1E, 0x40, 0x80)),
        ("⑦ PPT 制作 + 预答辩 + 修改", 6, 2, COLOR_PRIMARY),
    ]
    bar_top = 2.20
    bar_h = 0.34
    bar_gap = 0.06
    label_w = 2.50
    for i, (label, start, dur, color) in enumerate(tasks):
        y = bar_top + i * (bar_h + bar_gap)
        # Task label on left
        add_textbox(
            slide,
            left=0.55,
            top=y,
            width=label_w,
            height=bar_h,
            text=label,
            size=10,
            color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Bar
        bar_x = track_left + start * week_w + 0.05
        bar_w = dur * week_w - 0.10
        add_filled_rect(slide, bar_x, y + 0.05, bar_w, bar_h - 0.10, fill=color, corner=0.20)

    # Bottom: risk & mitigation
    add_textbox(
        slide,
        left=0.55,
        top=5.30,
        width=12.30,
        height=0.40,
        text="◆ 风险与保障措施",
        size=15,
        bold=True,
        color=COLOR_PRIMARY,
    )
    risk_left_w = 6.20
    add_filled_rect(slide, 0.55, 5.75, risk_left_w, 1.70, fill=COLOR_LIGHT_BG, corner=0.04)
    add_textbox(
        slide,
        left=0.65,
        top=5.80,
        width=risk_left_w,
        height=0.35,
        text="⚠️  潜在风险与应对",
        size=12,
        bold=True,
        color=COLOR_PRIMARY,
    )
    risks = [
        ("主观题评分准确性不达标", "→  语义相似度 + LLM 加权 + 5 组样本验证"),
        ("论文与代码冲突", "→  第 12 周完成代码，13-14 周专注论文"),
        ("DashScope API 不稳定", "→  预留本地模型接口（Qwen2.5 + Ollama）兜底"),
    ]
    box = slide.shapes.add_textbox(Inches(0.65), Inches(6.20), Inches(risk_left_w - 0.20), Inches(1.30))
    items = []
    for r, m in risks:
        items.append((f"• {r}", {"size": 11, "bold": True, "color": COLOR_TEXT, "space_after": 2}))
        items.append((f"  {m}", {"size": 10, "color": COLOR_ACCENT, "space_after": 6}))
    add_paragraphs(box.text_frame, items)

    safe_left = 6.95
    safe_w = 5.95
    add_filled_rect(slide, safe_left, 5.75, safe_w, 1.70, fill=COLOR_LIGHT_BG, corner=0.04)
    add_textbox(
        slide,
        left=safe_left + 0.10,
        top=5.80,
        width=safe_w,
        height=0.35,
        text="✓  保障措施",
        size=12,
        bold=True,
        color=COLOR_PRIMARY,
    )
    safes = [
        "✓ 已完成模块代码可继续打磨",
        "✓ 论文素材（架构图 / 截图 / 数据）已基本就绪",
        "✓ 关键技术难点已在前 8 周突破",
        "✓ 剩余多为业务层实现，工作量可控",
        "✓ 与导师每周 1 次进度沟通",
    ]
    box = slide.shapes.add_textbox(Inches(safe_left + 0.10), Inches(6.18), Inches(safe_w - 0.20), Inches(1.30))
    add_paragraphs(
        box.text_frame,
        [(s, {"size": 11, "color": COLOR_TEXT, "space_after": 4}) for s in safes],
    )

    add_page_number(slide, 12)


def build_p14_thanks(slide) -> None:
    """Page 14: 致谢. Operates on the template's existing thanks slide."""
    spTree = slide.shapes._spTree  # noqa: SLF001
    to_remove = []
    for shp in slide.shapes:
        st = shp.shape_type
        if st in (17, 14, 6):  # TEXT_BOX, PLACEHOLDER, GROUP
            to_remove.append(shp)
        elif st == 1:
            try:
                if shp.has_text_frame and shp.text_frame.text.strip():
                    to_remove.append(shp)
            except Exception:
                pass
    for shp in to_remove:
        sp = shp._element  # noqa: SLF001
        sp.getparent().remove(sp)

    # Big "谢谢" text
    add_textbox(
        slide,
        left=2.0,
        top=2.10,
        width=9.30,
        height=1.40,
        text="谢    谢",
        size=80,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # English thanks
    add_textbox(
        slide,
        left=2.0,
        top=3.55,
        width=9.30,
        height=0.5,
        text="THANK  YOU",
        size=18,
        color=COLOR_ACCENT,
        align=PP_ALIGN.CENTER,
        font_name=EN_FONT,
    )

    # Acknowledgments
    text = (
        "感谢指导教师 XXX 老师在选题方向、技术路线选择与项目推进中的悉心指导。\n"
        "感谢计算机科学与技术学院全体老师的悉心教诲。\n"
        "感谢同学们在开发过程中的讨论与协助，感谢家人对学业的长期支持。"
    )
    add_textbox(
        slide,
        left=1.5,
        top=4.70,
        width=10.3,
        height=1.20,
        text=text,
        size=14,
        color=COLOR_TEXT,
        align=PP_ALIGN.CENTER,
        line_spacing=1.5,
    )

    # Closing line
    add_textbox(
        slide,
        left=1.5,
        top=6.10,
        width=10.3,
        height=0.40,
        text="—— 请各位答辩老师批评指正！——",
        size=16,
        bold=True,
        color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide,
        left=1.5,
        top=6.65,
        width=10.3,
        height=0.40,
        text="李 鑫    ·    2026 年",
        size=12,
        color=COLOR_SUBTLE,
        align=PP_ALIGN.CENTER,
    )


# ------------------------------------------------------------------
# Main
# ------------------------------------------------------------------
def main() -> None:
    """Reuse the template's slides as canvases (13-page version):
    - Slide 0  -> Cover
    - Slides 1-11 -> 11 content pages (uses the template's content layouts)
    - Slide 26 -> Thanks
    Slides 12-25 are deleted as they aren't needed.
    """
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)
    print(f"Loaded template: {len(slides)} slides")

    # Build cover (template slide index 0)
    build_cover(slides[0])

    # Build 11 content pages on template slides 1..11
    builders = [
        build_p2_background,        # 选题背景
        build_p3_objective,         # 痛点 → 对策
        build_p4_use_case,          # 用户角色用例
        build_p5_architecture,      # 综合架构 + 技术栈
        build_p7_database,          # 数据库设计
        build_p8_teacher_overview,  # 教师侧总览
        build_p9_rag,               # 模块 1 RAG
        build_p10_lesson_plan,      # 模块 2 备课
        build_p11_exercise,         # 模块 3 出题
        build_p12_progress,         # 进度对照
        build_p13_remaining,        # 剩余工作
    ]
    for i, builder in enumerate(builders):
        slide_idx = 1 + i  # template slides 1..11
        builder(slides[slide_idx])

    # Build thanks page on template slide 26
    build_p14_thanks(slides[26])

    # Delete unused template slides (indices 12..25 in original order).
    # Delete from highest index to avoid reindex problems.
    for idx in range(25, 11, -1):
        delete_slide_at(prs, idx)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
