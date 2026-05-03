"""Insert / rebuild slide 11 — Agent 编排层 (LangGraph + MCP + Skills).

This script touches ONLY the new page. Every other slide (including
slides 1-10 and the original "进度对照 / 剩余工作 / 致谢" trailer
that gets shifted to slides 12-14) is left exactly as-is — we never
clear, mutate, or re-style their existing shapes / pictures.

Idempotent:
  - If the deck currently has 13 slides (pre-insertion state), we
    duplicate slide 10 as a blank canvas and insert it at index 10
    (1-based slide 11), then rebuild it.
  - If the deck already has 14 slides (post-insertion state), we
    simply rebuild slide 11 in place.

Run:
    .ppt_venv/bin/python scripts/add_p11_agent_orchestration.py
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent))
from build_defense_pptx import (  # noqa: E402
    COLOR_ACCENT,
    COLOR_LIGHT_BG,
    COLOR_PRIMARY,
    COLOR_SUBTLE,
    COLOR_TEXT,
    EN_FONT,
    add_filled_rect,
    add_textbox,
    add_title_bar,
)
from update_defense_pptx import (  # noqa: E402
    clear_slide_keep_logo_only,
    duplicate_slide_after,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

PPT_PATH = Path("/Users/daniel/2026design/App/中期答辩PPT.pptx")

# Layer accent colors (consistent with existing architecture page):
COLOR_LANGGRAPH = RGBColor(0x1E, 0x40, 0x80)   # deep blue – brain
COLOR_SKILLS = RGBColor(0xD9, 0x82, 0x2E)      # orange – workflow
COLOR_MCP = RGBColor(0x2E, 0x9C, 0x6E)         # green – atomic tools
COLOR_HIGHLIGHT = RGBColor(0xC0, 0x55, 0x55)   # red highlight


# =============================================================================
# Slide builder
# =============================================================================


def rebuild_p11_agent(slide) -> None:
    """Build the "Agent 编排层 — LangGraph + MCP + Skills" slide.

    Layout (13.333" × 7.5", title bar leaves usable area 1.10 → 7.40):
       1.10–1.46  "承接第 4 页" banner (one line, light)
       1.55–3.55  Three-layer cards (horizontal: MCP → Skills → LangGraph)
       3.62–3.85  "类比" caption row (一行金句：食材 / 菜谱 / 主厨)
       3.95–6.05  Left  LangGraph 5-node state machine (w≈5.30)
       3.95–6.05  Right 多智能体协作 5 cards (w≈7.00, single row of 5)
       6.13–6.95  MCP 5 tools chip row + 3 Skill chips
       7.00–7.35  Bottom golden banner
    """
    clear_slide_keep_logo_only(slide)
    add_title_bar(slide, "Agent 编排层 — LangGraph 多智能体 + MCP + Skills")

    # ====================================================================
    # 1. 承接说明 banner
    # ====================================================================
    add_filled_rect(slide, 0.55, 1.10, 12.30, 0.36,
                    fill=COLOR_LIGHT_BG, corner=0.04)
    add_textbox(
        slide, left=0.55, top=1.10, width=12.30, height=0.36,
        text="◆ 承接第 4 页架构图「⑤ Agent 编排层」放大展开 — 把 LLM 不确定性转化为工程确定性的核心引擎",
        size=12, bold=True, color=COLOR_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ====================================================================
    # 2. 三层暴露：3 张横向卡片
    #    顺序（自底向上）：① MCP Server → ② Agent Skills → ③ LangGraph
    # ====================================================================
    cards_top = 1.55
    cards_h = 2.00
    card_gap = 0.18
    n_cards = 3
    card_w = (12.30 - (n_cards - 1) * card_gap) / n_cards  # ≈ 3.98

    layers = [
        {
            "tag": "①",
            "name": "MCP Server",
            "role": "原子能力 · 协议级",
            "color": COLOR_MCP,
            "metaphor": "食材",
            "highlights": [
                ("5 个原子 Tool · 标准协议",
                 "search_kb / lesson_outline / generate_exercise\n"
                 "grade_answer / get_mastery\n"
                 "Claude Desktop · Cursor 等客户端可直接调用"),
                ("无状态 · 单一职责",
                 "每个 Tool 输入输出 JSON Schema 严格约束"),
            ],
        },
        {
            "tag": "②",
            "name": "Agent Skills",
            "role": "工作流 · 菜谱级",
            "color": COLOR_SKILLS,
            "metaphor": "菜谱",
            "highlights": [
                ("3 个高层 Skill · 场景化封装",
                 "prepare-class（教师一键备课）\n"
                 "personalized-practice（学生个性化练习）\n"
                 "grade-essay（主观题批量批改）"),
                ("串联多个 MCP Tool",
                 "封装常用业务流程，对外暴露单一入口"),
            ],
        },
        {
            "tag": "③",
            "name": "LangGraph",
            "role": "决策大脑 · 智能体编排",
            "color": COLOR_LANGGRAPH,
            "metaphor": "主厨",
            "highlights": [
                ("5 节点状态机 + 反思回边",
                 "意图识别 → 任务拆解 → 工具调用 →\n"
                 "反思纠错 ─失败重试─↑ → 输出聚合\n"
                 "Schema 校验 + LLM-as-Judge 自愈"),
                ("多 Agent 异步协作",
                 "检索 / 答疑 / 学情 / 推荐 / 反思 5 个子 Agent"),
            ],
        },
    ]

    for ci, layer in enumerate(layers):
        x = 0.55 + ci * (card_w + card_gap)
        # Card body
        add_filled_rect(slide, x, cards_top, card_w, cards_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        # Header strip (color)
        head_h = 0.55
        add_filled_rect(slide, x, cards_top, card_w, head_h,
                        fill=layer["color"], corner=0.06)
        # Tag chip in header (left)
        add_textbox(
            slide, left=x + 0.10, top=cards_top + 0.04,
            width=0.50, height=head_h - 0.08,
            text=layer["tag"], size=22, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        # Layer name (header center)
        add_textbox(
            slide, left=x + 0.65, top=cards_top + 0.04,
            width=card_w - 1.45, height=0.30,
            text=layer["name"], size=15, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            anchor=MSO_ANCHOR.MIDDLE, font_name=EN_FONT,
        )
        add_textbox(
            slide, left=x + 0.65, top=cards_top + 0.30,
            width=card_w - 1.45, height=0.22,
            text=layer["role"], size=10,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Metaphor pill on the far right of header
        meta_w = 0.70
        meta_pill_x = x + card_w - meta_w - 0.10
        add_filled_rect(
            slide, meta_pill_x, cards_top + 0.13,
            meta_w, head_h - 0.26,
            fill=RGBColor(0xFF, 0xFF, 0xFF), corner=0.20,
        )
        add_textbox(
            slide, left=meta_pill_x, top=cards_top + 0.13,
            width=meta_w, height=head_h - 0.26,
            text=layer["metaphor"], size=11, bold=True,
            color=layer["color"],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Body: highlight rows with generous spacing
        body_top = cards_top + head_h + 0.08
        body_h = cards_h - head_h - 0.16
        n_rows = len(layer["highlights"])
        row_gap = 0.08
        row_h = (body_h - (n_rows - 1) * row_gap) / n_rows
        title_h = 0.22
        for hi, (h_title, h_desc) in enumerate(layer["highlights"]):
            ry = body_top + hi * (row_h + row_gap)
            # tiny color accent bar on the left
            add_filled_rect(
                slide, x + 0.10, ry + 0.04, 0.08, row_h - 0.08,
                fill=layer["color"], corner=0.20,
            )
            add_textbox(
                slide, left=x + 0.25, top=ry,
                width=card_w - 0.35, height=title_h,
                text=h_title, size=10, bold=True,
                color=layer["color"],
                anchor=MSO_ANCHOR.TOP,
            )
            add_textbox(
                slide, left=x + 0.25, top=ry + title_h,
                width=card_w - 0.35, height=row_h - title_h,
                text=h_desc, size=9, color=COLOR_TEXT,
                line_spacing=1.20, anchor=MSO_ANCHOR.TOP,
            )

    # ====================================================================
    # 3. 类比金句（一行）
    # ====================================================================
    add_textbox(
        slide, left=0.55, top=cards_top + cards_h + 0.07,
        width=12.30, height=0.22,
        text="★ 类比：MCP = 食材  ·  Skills = 菜谱  ·  LangGraph = 主厨    "
             "— 三层职责清晰、各司其职，命中工业界主流 Agent 范式",
        size=11, bold=True, color=COLOR_HIGHLIGHT,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # ====================================================================
    # 4. 中下区域：左 LangGraph 状态机 | 右 多智能体协作
    # ====================================================================
    mid_top = 3.95
    mid_h = 2.10
    left_w = 5.30
    left_x = 0.55
    right_x = 0.55 + left_w + 0.20
    right_w = 12.30 - left_w - 0.20  # ≈ 6.80

    # ---- 左：LangGraph 5 节点状态机（垂直 pipeline）----
    add_filled_rect(slide, left_x, mid_top, left_w, 0.34,
                    fill=COLOR_LANGGRAPH, corner=0.04)
    add_textbox(
        slide, left=left_x, top=mid_top, width=left_w, height=0.34,
        text="◆ LangGraph 状态机（5 节点 + 反思回边）",
        size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    nodes = [
        ("①", "意图识别", "Qwen-Plus 判定任务类型"),
        ("②", "任务拆解", "把高层目标拆成原子步骤 DAG"),
        ("③", "工具调用", "Function Calling → MCP / Skill / 子 Agent"),
        ("④", "反思纠错", "Schema 校验 + LLM-as-Judge ─失败回边─↑"),
        ("⑤", "输出聚合", "SSE 流式推送前端"),
    ]
    body_top = mid_top + 0.40
    body_h = mid_h - 0.42
    n_nodes = len(nodes)
    node_gap = 0.04
    node_h = (body_h - (n_nodes - 1) * node_gap) / n_nodes
    for i, (no, name, desc) in enumerate(nodes):
        y = body_top + i * (node_h + node_gap)
        add_filled_rect(slide, left_x, y, left_w, node_h,
                        fill=COLOR_LIGHT_BG, corner=0.04)
        # Number badge
        badge_w = 0.42
        # node ④ uses HIGHLIGHT color to emphasise the reflection回边
        node_color = COLOR_HIGHLIGHT if i == 3 else COLOR_LANGGRAPH
        add_filled_rect(slide, left_x, y, badge_w, node_h,
                        fill=node_color, corner=0.04)
        add_textbox(
            slide, left=left_x, top=y, width=badge_w, height=node_h,
            text=no, size=14, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left=left_x + badge_w + 0.10, top=y + 0.02,
            width=1.30, height=node_h - 0.04,
            text=name, size=11, bold=True, color=node_color,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left=left_x + badge_w + 1.45, top=y + 0.02,
            width=left_w - badge_w - 1.55, height=node_h - 0.04,
            text=desc, size=9, color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15,
        )
        if i < n_nodes - 1:
            arrow_y = y + node_h + node_gap * 0.05
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                Inches(left_x + badge_w / 2 - 0.05),
                Inches(arrow_y),
                Inches(0.10), Inches(node_gap * 0.90),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_SUBTLE
            arrow.line.fill.background()

    # ---- 右：多智能体协作 5 个 Agent 卡片（单行 5 列）----
    add_filled_rect(slide, right_x, mid_top, right_w, 0.34,
                    fill=COLOR_SKILLS, corner=0.04)
    add_textbox(
        slide, left=right_x, top=mid_top, width=right_w, height=0.34,
        text="◆ 多智能体协作 · 学生提问 → 个性化推荐闭环",
        size=11, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    agents = [
        ("检索", "Retrieve", RGBColor(0x4A, 0x90, 0xE2),
         "search_kb\nBM25 + 向量\n+ Rerank"),
        ("答疑", "Answer", RGBColor(0x2E, 0x6F, 0xB5),
         "Qwen-Plus 流式\n+ 引用回填\n[1] [2]"),
        ("学情", "Mastery", RGBColor(0xE5, 0xA5, 0x3B),
         "get_mastery\nEMA 平滑\n后向量"),
        ("推荐", "Recommend", RGBColor(0x2E, 0x9C, 0x6E),
         "generate\n_exercise\n按薄弱知识点"),
        ("反思", "Reflect", COLOR_HIGHLIGHT,
         "LLM-as-Judge\n+ JSON\nSchema"),
    ]
    a_top = mid_top + 0.40
    a_h = mid_h - 0.42
    n_agents = len(agents)
    a_gap = 0.10
    a_w = (right_w - (n_agents - 1) * a_gap) / n_agents  # ≈ 1.28
    for i, (name, en, color, desc) in enumerate(agents):
        ax = right_x + i * (a_w + a_gap)
        add_filled_rect(slide, ax, a_top, a_w, a_h,
                        fill=COLOR_LIGHT_BG, corner=0.06)
        # color header strip
        head_h2 = 0.46
        add_filled_rect(slide, ax, a_top, a_w, head_h2,
                        fill=color, corner=0.06)
        add_textbox(
            slide, left=ax, top=a_top + 0.02, width=a_w, height=0.24,
            text=name, size=12, bold=True,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER,
        )
        add_textbox(
            slide, left=ax, top=a_top + 0.24, width=a_w, height=0.20,
            text=en + " Agent", size=8,
            color=RGBColor(0xFF, 0xFF, 0xFF),
            align=PP_ALIGN.CENTER, font_name=EN_FONT,
        )
        # Body: tool / desc
        add_textbox(
            slide, left=ax + 0.06, top=a_top + head_h2 + 0.06,
            width=a_w - 0.12, height=a_h - head_h2 - 0.12,
            text=desc, size=9, color=COLOR_TEXT,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            line_spacing=1.25,
        )
        # Down-arrow connecting consecutive agents
        if i < n_agents - 1:
            arr_x = ax + a_w + 0.005
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arr_x), Inches(a_top + a_h / 2 - 0.08),
                Inches(a_gap - 0.02), Inches(0.16),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_PRIMARY
            arrow.line.fill.background()

    # ====================================================================
    # 5. 下区域：MCP 5 个工具 chip 行 + Skills 3 chip 行
    # ====================================================================
    tools_top = 6.13
    tools_h = 0.40
    add_textbox(
        slide, left=0.55, top=tools_top, width=2.00, height=tools_h,
        text="MCP Tools ·",
        size=11, bold=True, color=COLOR_MCP,
        anchor=MSO_ANCHOR.MIDDLE, font_name=EN_FONT,
    )
    tools = [
        ("search_kb", "知识库混合检索"),
        ("lesson_outline", "教师备课"),
        ("generate_exercise", "智能出题"),
        ("grade_answer", "主观题评分"),
        ("get_mastery", "学情分析"),
    ]
    chip_left0 = 2.45
    chip_total_w = 12.85 - chip_left0
    chip_gap = 0.10
    chip_w = (chip_total_w - (len(tools) - 1) * chip_gap) / len(tools)  # ≈ 1.98
    for i, (tname, tdesc) in enumerate(tools):
        cx = chip_left0 + i * (chip_w + chip_gap)
        add_filled_rect(slide, cx, tools_top, chip_w, tools_h,
                        fill=COLOR_LIGHT_BG, corner=0.10)
        add_filled_rect(slide, cx, tools_top, 0.10, tools_h,
                        fill=COLOR_MCP, corner=0.10)
        add_textbox(
            slide, left=cx + 0.15, top=tools_top + 0.02,
            width=chip_w - 0.20, height=0.20,
            text=tname, size=9, bold=True, color=COLOR_MCP,
            font_name=EN_FONT, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_textbox(
            slide, left=cx + 0.15, top=tools_top + 0.20,
            width=chip_w - 0.20, height=0.18,
            text=tdesc, size=8, color=COLOR_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    skills_top = tools_top + tools_h + 0.08
    skills_h = 0.32
    add_textbox(
        slide, left=0.55, top=skills_top, width=2.00, height=skills_h,
        text="Agent Skills ·",
        size=11, bold=True, color=COLOR_SKILLS,
        anchor=MSO_ANCHOR.MIDDLE, font_name=EN_FONT,
    )
    skills = [
        ("prepare-class", "search_kb → lesson_outline"),
        ("personalized-practice", "get_mastery → generate_exercise"),
        ("grade-essay", "grade_answer ×N → 聚合反馈"),
    ]
    skill_left0 = 2.45
    skill_total_w = 12.85 - skill_left0
    skill_gap = 0.12
    skill_w = (skill_total_w - (len(skills) - 1) * skill_gap) / len(skills)
    for i, (sname, sdesc) in enumerate(skills):
        sx = skill_left0 + i * (skill_w + skill_gap)
        add_filled_rect(slide, sx, skills_top, skill_w, skills_h,
                        fill=COLOR_LIGHT_BG, corner=0.10)
        add_filled_rect(slide, sx, skills_top, 0.10, skills_h,
                        fill=COLOR_SKILLS, corner=0.10)
        # Skill name (left)
        name_w = 1.55
        add_textbox(
            slide, left=sx + 0.15, top=skills_top, width=name_w,
            height=skills_h, text=sname, size=10, bold=True,
            color=COLOR_SKILLS, font_name=EN_FONT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Tool chain (right)
        add_textbox(
            slide, left=sx + 0.15 + name_w, top=skills_top,
            width=skill_w - 0.20 - name_w, height=skills_h,
            text=sdesc, size=9, color=COLOR_TEXT, font_name=EN_FONT,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    # ====================================================================
    # 6. 底部金句横幅
    # ====================================================================
    banner_top = skills_top + skills_h + 0.08  # ≈ 7.01
    add_filled_rect(slide, 0.55, banner_top, 12.30, 0.30,
                    fill=COLOR_PRIMARY,
                    shape_type=MSO_SHAPE.ROUNDED_RECTANGLE,
                    corner=0.40)
    add_textbox(
        slide, left=0.55, top=banner_top, width=12.30, height=0.30,
        text="MCP 标准化原子能力 · Skills 封装常用工作流 · LangGraph 做智能决策 — "
             "把 LLM 不确定性转化为工程确定性",
        size=12, bold=True, color=RGBColor(0xFF, 0xCC, 0x00),
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )


# =============================================================================
# Main
# =============================================================================


def main() -> int:
    if not PPT_PATH.exists():
        print(f"[ERR] PPT not found: {PPT_PATH}")
        return 1

    prs = Presentation(PPT_PATH)
    n_before = len(prs.slides)
    print(f"Loaded: {PPT_PATH.name}  ({n_before} slides)")

    # ----------------------------------------------------------------
    # Step 1 — Insert (or reuse) slide 11 right after slide 10.
    # ----------------------------------------------------------------
    if n_before == 13:
        # Pre-insertion state: duplicate slide 10 as a blank container,
        # which becomes the new slide 11.
        new_slide = duplicate_slide_after(prs, source_idx=9)
        print("  [insert ] duplicated slide 10 → new container at slide 11")
    elif n_before == 14:
        # Already inserted previously — just rebuild slide 11 in place.
        new_slide = list(prs.slides)[10]
        print("  [skip   ] deck already 14 slides; reusing existing slide 11")
    else:
        print(f"  [WARN   ] unexpected deck size {n_before}; aborting")
        return 1

    # ----------------------------------------------------------------
    # Step 2 — Build the new "Agent 编排层" page on slide 11.
    # ----------------------------------------------------------------
    rebuild_p11_agent(new_slide)
    print("  [rebuilt] slide 11 (Agent 编排层 — LangGraph + MCP + Skills)")

    prs.save(PPT_PATH)
    print(f"Saved: {PPT_PATH.name}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    sys.exit(main())
