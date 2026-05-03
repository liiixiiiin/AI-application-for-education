"""Inspect the PPT template structure: slides, layouts, masters."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

TEMPLATE_PATH = Path(
    "/Users/daniel/答辩PPT/中国石油大学（华东）PPT模板/学海遨游，逐梦石大-学术系列_副本.pptx"
)


def main() -> None:
    prs = Presentation(TEMPLATE_PATH)
    print(f"Slide width: {prs.slide_width} EMU = {Emu(prs.slide_width).inches:.2f} in")
    print(f"Slide height: {prs.slide_height} EMU = {Emu(prs.slide_height).inches:.2f} in")
    print(f"Total slides: {len(prs.slides)}")
    print(f"Slide masters: {len(prs.slide_masters)}")
    for mi, master in enumerate(prs.slide_masters):
        print(f"\n=== Master {mi} ===")
        print(f"  Layouts: {len(master.slide_layouts)}")
        for li, layout in enumerate(master.slide_layouts):
            placeholders = [
                f"idx={ph.placeholder_format.idx} type={ph.placeholder_format.type} name={ph.name!r}"
                for ph in layout.placeholders
            ]
            print(f"    Layout {li}: name={layout.name!r}")
            for p in placeholders:
                print(f"      {p}")

    print("\n=== Slides ===")
    for si, slide in enumerate(prs.slides):
        layout_name = slide.slide_layout.name
        n_shapes = len(slide.shapes)
        print(f"  Slide {si + 1}: layout={layout_name!r} shapes={n_shapes}")


if __name__ == "__main__":
    main()
